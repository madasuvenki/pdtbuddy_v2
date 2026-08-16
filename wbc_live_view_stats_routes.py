import hashlib
import json
import os
import re
from datetime import date, datetime, timedelta
from glob import glob
from typing import Any, Dict, List, Optional, Tuple


import io
from flask import Blueprint, jsonify, render_template, request, send_file
from flask_login import current_user, login_required

from config import ADMIN_USERS, BU_DATABASE_MAPPING, JIRA_PDT_FILTER_ID, TARGET_GROUP, VIEWER_OVERRIDE_USERS
from dashboard_common import get_mysql_connection_db
from live_view_stats_routes import _sheet_to_payload, _workbook_sheets


wbc_live_view_stats_bp = Blueprint("wbc_live_view_stats_bp", __name__)

_WBC_ROOT = os.environ.get("WBC_LIVE_VIEW_STATS_ROOT", r"C:\Dropbox\WBC_Scrum_DB")
_WBC_DB_FILES = os.environ.get("WBC_LIVE_VIEW_STATS_DB_FILES", os.path.join(_WBC_ROOT, "DB_Files"))
_WBC_FR_FILES = os.environ.get(
    "WBC_LIVE_VIEW_FR_FILES",
    r"\\sphere\pdtqipl_internal\PDTBuddy\live_status_publish\WBC\FRs",
)
_DATA_ROOT = os.environ.get("PDTBUDDY_DATA_ROOT", r"\\Sphere\pdtqipl_internal\PDTBuddy")
_LOCAL_ROOT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data", "wbc_live_view_stats")
_WBC_SCHEMA = str(BU_DATABASE_MAPPING.get("WBC") or "pdt_stats_wbc").strip("`")


def _can_edit() -> bool:
    uid = str(getattr(current_user, "id", "") or "").strip().lower()
    if not uid or uid in VIEWER_OVERRIDE_USERS:
        return False
    if uid in ADMIN_USERS:
        return True
    try:
        import app as _app
        return bool(_app.is_user_in_group(uid, TARGET_GROUP))
    except Exception:
        return False


def _slug(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9_.-]+", "_", str(value or "").strip()).strip("_") or "target"


def _store_dir(create: bool = True) -> str:
    folder = os.path.join(_DATA_ROOT, "managed_excel", "WBC", "LIVE_VIEW_STATS")
    if create:
        try:
            os.makedirs(folder, exist_ok=True)
            probe = os.path.join(folder, ".write_probe")
            with open(probe, "w", encoding="utf-8") as fh:
                fh.write("ok")
            os.remove(probe)
            return folder
        except Exception:
            pass
        os.makedirs(_LOCAL_ROOT, exist_ok=True)
    return folder if os.path.isdir(folder) else _LOCAL_ROOT


def _config_path() -> str:
    return os.path.join(_store_dir(), "config.json")


def _target_json_path(target_key: str) -> str:
    return os.path.join(_store_dir(), f"target_{_slug(target_key)}.json")


def _mtbf_json_path(target_key: str) -> str:
    return os.path.join(_store_dir(), f"mtbf_{_slug(target_key)}_Mainline_Build_Details.json")


def _dashboard_mtbf_target_name(target_key: str, target: dict = None) -> str:
    """Derive the dashboard target name from the WBC target.

    The dashboard stores target names in dashboard_status.target_name using the
    format: label.replace(".", "_").lower()
    e.g. WBC label "Kobuk.LE.1.1" → dashboard target name "kobuk_le_1_1"

    This must match the path used by dashboard_routes._mtbf_json_dir().
    """
    label = str((target or {}).get("label") or "").strip()
    if label:
        # Convert WBC label to dashboard target name format
        # "Kobuk.LE.1.1" → "kobuk_le_1_1"
        return label.replace(".", "_").lower()
    # Fallback: use the WBC key as-is
    return target_key


def _dashboard_mtbf_json_path(target_key: str, target: dict = None) -> str:
    """Return the internal dashboard MTBF JSON path for this WBC target.

    Uses the dashboard target name format (label.replace(".", "_").lower()) so
    the path matches what dashboard_routes._load_mtbf_json_payload reads from:
      PDTBUDDY_DATA_ROOT/managed_excel/WBC/{dashboard_target_name}/mtbf_mtbf.json

    e.g. WBC key "Kobuk11" / label "Kobuk.LE.1.1"
         → managed_excel/WBC/kobuk_le_1_1/mtbf_mtbf.json
    """
    import re as _re
    def _safe_slug(v):
        return _re.sub(r'[^A-Za-z0-9_.-]+', '_', str(v or '').strip()).strip('._') or 'target'
    dash_name = _dashboard_mtbf_target_name(target_key, target)
    return os.path.join(_DATA_ROOT, "managed_excel", "WBC", _safe_slug(dash_name), "mtbf_mtbf.json")


def _sync_to_dashboard_mtbf_json(target_key: str, data: dict, target: dict = None) -> None:
    """Write WBC MTBF chart_rows to the internal dashboard MTBF JSON path.

    Converts the WBC chart_rows format to the internal rows format so that
    /api/dashboard/<target>/excel/full_table returns the same data as the
    WBC Live View Stats page.

    The path is derived from the WBC label (e.g. "Kobuk.LE.1.1" → "kobuk_le_1_1")
    to match the dashboard_status.target_name format used by dashboard_routes.
    """
    try:
        chart_rows = data.get("chart_rows") or []
        if not chart_rows:
            return
        internal_rows = []
        for i, cr in enumerate(chart_rows, 1):
            build = str(cr.get("crm_build_id") or cr.get("meta_id") or cr.get("build") or "").strip()
            if not build:
                continue
            hours   = float(cr.get("hours") or 0)
            crashes = int(cr.get("total_crashes") or cr.get("crash") or 0)
            mtbf    = float(cr.get("mtbf") or 0)
            if not mtbf and hours and crashes:
                mtbf = round(hours / crashes, 2)
            internal_rows.append({
                "id":            str(cr.get("id") or f"wbc_{i}"),
                "meta_id":       build,
                "build":         build,
                "build_full":    build,
                "date":          str(cr.get("date") or "")[:10],
                "hours":         round(hours, 2),
                "total_crashes": crashes,
                "mtbf":          round(mtbf, 2),
                "comments":      "",
                "_source":       "wbc_live_view_stats",
            })
        if not internal_rows:
            return
        dash_name = _dashboard_mtbf_target_name(target_key, target)
        payload = {
            "target":     dash_name,
            "view":       "MTBF",
            "headers":    ["Meta ID", "Build(s)", "Date", "Hours", "Total Crashes", "MTBF", "Comments"],
            "rows":       internal_rows,
            "updated_at": data.get("updated_at") or "",
            "_source":    "wbc_live_view_stats",
        }
        path = _dashboard_mtbf_json_path(target_key, target)
        os.makedirs(os.path.dirname(path), exist_ok=True)
        tmp = path + ".tmp"
        import json as _json
        with open(tmp, "w", encoding="utf-8") as fh:
            _json.dump(payload, fh, indent=2, ensure_ascii=False)
        os.replace(tmp, path)
    except Exception as _e:
        import logging as _log
        _log.getLogger(__name__).debug("[WBC MTBF SYNC] failed for %s: %s", target_key, _e)


def _mtbf_aux_dir() -> str:
    folder = os.path.join(_store_dir(), "mtbf")
    os.makedirs(folder, exist_ok=True)
    return folder


def _fr_json_dir() -> str:
    folder = os.path.join(_store_dir(), "fr_analysis")
    os.makedirs(folder, exist_ok=True)
    return folder


def _fr_json_path(target_key: str) -> str:
    return os.path.join(_fr_json_dir(), f"fr_analysis_{_slug(target_key)}.json")


def _overview_summary_path(target_key: str) -> str:
    return os.path.join(_mtbf_aux_dir(), f"overview_summary_{_slug(target_key)}.json")


def _running_report_cache_path(target_key: str) -> str:
    folder = os.path.join(_store_dir(), "running_build_reports")
    os.makedirs(folder, exist_ok=True)
    return os.path.join(folder, f"running_build_report_{_slug(target_key)}.json")


def _running_build_jql_cache_path(target_key: str, build_id: str) -> str:
    folder = os.path.join(_store_dir(), "running_build_reports", "jql_builds", _slug(target_key))
    os.makedirs(folder, exist_ok=True)
    key = hashlib.sha1(_build_tail(build_id).encode("utf-8", errors="ignore")).hexdigest()[:18]
    return os.path.join(folder, f"{key}.json")


def _wbc_saved_jql_domain(target: Dict[str, str] | str | None = None) -> str:
    """Use a distinct saved-JQL namespace for every WBC PL/target."""
    if isinstance(target, dict):
        target = target.get("key") or target.get("name") or target.get("label")
    key = _slug(str(target or "WBC"))
    return f"WBC_{key}"


def _wbc_saved_jql_filter_id(value: Any) -> str:
    """Extract JIRA saved-filter ID from ID, filter=ID JQL, or JIRA filter URL."""
    text = str(value or "").strip()
    if not text:
        return ""
    if text.isdigit():
        return text
    try:
        from urllib.parse import parse_qs, urlparse
        qs = parse_qs(urlparse(text).query or "")
        for key in ("filter", "filterId"):
            val = str((qs.get(key) or [""])[0]).strip()
            if val.isdigit():
                return val
    except Exception:
        pass
    match = re.match(r"^\s*filter(?:Id)?\s*=\s*(\d+)\s*(?:ORDER\s+BY\s+.+)?$", text, flags=re.I)
    if match:
        return match.group(1)
    match = re.search(r"[?&]filter(?:Id)?=(\d+)", text, flags=re.I)
    return match.group(1) if match else ""


def _wbc_resolve_jira_filter_jql(filter_id: str) -> str:
    """Resolve a saved JIRA filter ID inside WBC only, so WBC never uses stale saved-filter text."""
    filter_id = str(filter_id or "").strip()
    if not filter_id:
        return ""
    import sys as _sys
    scripts_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "scripts")
    if scripts_dir not in _sys.path:
        _sys.path.insert(0, scripts_dir)
    from config import JIRA_PASSWORD, JIRA_SERVER_ENDPOINT, JIRA_USER
    from fetch_consolidated_report import connect_jira
    jira_obj = connect_jira(JIRA_USER, JIRA_PASSWORD, JIRA_SERVER_ENDPOINT)
    filt = jira_obj.filter(filter_id)
    return str(getattr(filt, "jql", "") or "").strip()


def _wbc_resolve_saved_jql(raw_jql: Any) -> Tuple[str, str, bool, str]:
    """Return (effective_jql, filter_id, resolved, error)."""
    raw = str(raw_jql or "").strip()
    filter_id = _wbc_saved_jql_filter_id(raw)
    if not filter_id:
        return raw, "", False, ""
    try:
        latest = _wbc_resolve_jira_filter_jql(filter_id)
        if latest:
            return latest, filter_id, True, ""
        return raw, filter_id, False, "Filter lookup returned empty JQL"
    except Exception as exc:
        return raw, filter_id, False, str(exc)


def _wbc_extract_build_id_from_jql(value: Any) -> str:
    text = str(value or "")
    patterns = [
        r"\b[A-Z][A-Z0-9_.]*\.LE\.[0-9.]+-[0-9]{3,6}-[A-Z0-9_.-]+(?:-[0-9]+)?\b",
        r"\b[A-Z][A-Z0-9_.-]+-[0-9]{3,6}-[A-Z0-9_.-]+(?:-[0-9]+)?\b",
        r"\b(?:META|BUILD)[-_ ]?0*([0-9]{3,6})\b",
    ]
    for pattern in patterns:
        match = re.search(pattern, text, flags=re.I)
        if match:
            return match.group(0)
    quoted = re.findall(r'"([^"\r\n]{6,160})"', text)
    return next((q for q in quoted if re.search(r"\d{3,6}", q) and re.search(r"[A-Za-z]", q)), "")





def _json_default(value: Any) -> str:
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    return str(value)


def _write_json(path: str, payload: Any) -> None:
    import time as _time
    os.makedirs(os.path.dirname(path), exist_ok=True)
    tmp = path + ".tmp"
    with open(tmp, "w", encoding="utf-8") as fh:
        json.dump(payload, fh, indent=2, ensure_ascii=False, default=_json_default)
    # os.replace() can fail on Windows network shares (WinError 5 / WinError 32)
    # when the destination is locked by antivirus, indexer, or a share-level lock.
    # Strategy: retry up to 3x with back-off → delete-then-rename → direct overwrite.
    last_err = None
    for attempt in range(3):
        try:
            os.replace(tmp, path)
            last_err = None
            break
        except OSError as exc:
            last_err = exc
            _time.sleep(0.05 * (attempt + 1))   # 50 ms, 100 ms, 150 ms
    if last_err is not None:
        logger.warning("[wbc] os.replace failed (%s); trying delete+rename", last_err)
        try:
            if os.path.exists(path):
                os.remove(path)
            os.rename(tmp, path)
        except OSError as exc2:
            logger.warning("[wbc] delete+rename failed (%s); direct-write fallback", exc2)
            try:
                with open(path, "w", encoding="utf-8") as fh:
                    json.dump(payload, fh, indent=2, ensure_ascii=False, default=_json_default)
            finally:
                try:
                    os.remove(tmp)
                except OSError:
                    pass


def _read_json(path: str, default: Any) -> Any:
    try:
        with open(path, "r", encoding="utf-8") as fh:
            return json.load(fh)
    except Exception:
        return default



# ---------------------------------------------------------------------------
# Friendly label mapping: raw Excel filename stem -> display label
# ---------------------------------------------------------------------------
_LABEL_OVERRIDES = {
    "Kobuk11": "Kobuk.LE.1.1", "Kobuk31": "Kobuk.LE.3.1",
    "Kobuk_LE11": "Kobuk.LE.1.1", "Kobuk_LE31": "Kobuk.LE.3.1",
    "Kobuk_Device": "Kobuk.LE.1.0",
    "Pinnacles_1_0": "Pinnacles.LE.1.0", "Pinnacles_1_2": "Pinnacles.LE.1.2",
    "Pinnacles_2_0": "Pinnacles.LE.2.0", "Pinnacles_2_2": "Pinnacles.LE.2.2",
    "Pinnacles_2_3": "Pinnacles.LE.2.3",
    "Pinnacles.1.0": "Pinnacles.LE.1.0", "Pinnacles.1.2": "Pinnacles.LE.1.2",
    "Pinnacles.2.0": "Pinnacles.LE.2.0", "Pinnacles.2.2": "Pinnacles.LE.2.2",
    "Pinnacles.2.3": "Pinnacles.LE.2.3",
    "Kuno_LE11": "Kuno.LE.1.1", "Kuno_LE_1_1": "Kuno.LE.1.1",
    "Kuno_LE_1_0": "Kuno.LE.1.0", "Kuno_LE": "Kuno.LE.1.0",
    "Kuno_TX": "Kuno.TX.1.0",
    "QMB415_LE": "QMB415.LE.1.0", "QMB415": "QMB415.LA.1.0", "QMB715": "QMB715.LA.1.0",
    "Tarang10": "Tarang.LE.1.0", "Tarang_LE_1_0": "Tarang.LE.1.0",
}


def _friendly_label(stem: str) -> str:
    """Convert raw Excel filename stem to human-friendly project label."""
    for k, v in _LABEL_OVERRIDES.items():
        if k.lower() == stem.lower():
            return v
    m = re.match(r"(?i)^([A-Za-z]+)(\d)(\d)$", stem)
    if m: return f"{m.group(1).capitalize()}.LE.{m.group(2)}.{m.group(3)}"
    m = re.match(r"(?i)^([A-Za-z]+)LE(\d)(\d)$", stem)
    if m: return f"{m.group(1).capitalize()}.LE.{m.group(2)}.{m.group(3)}"
    m = re.match(r"(?i)^([A-Za-z]+)[_\s]LE(\d)(\d)$", stem)
    if m: return f"{m.group(1).capitalize()}.LE.{m.group(2)}.{m.group(3)}"
    m = re.match(r"(?i)^([A-Za-z]+)[_\s](\d)[_\s](\d)$", stem)
    if m: return f"{m.group(1).capitalize()}.LE.{m.group(2)}.{m.group(3)}"
    return stem.replace("_", ".")


def _target_from_excel(path: str) -> Dict[str, str]:
    base = os.path.basename(path)
    name = re.sub(r"(?i)_Device_Deployment\.(xlsx|xlsm)$", "", base).strip()
    label = _friendly_label(name)
    return {"key": _slug(name), "name": name, "label": label, "excel_path": path}

def _find_target_excel(target: Dict[str, str], db_cfg: Dict[str, str] = None) -> str:
    explicit = str((db_cfg or {}).get("excel_path") or target.get("excel_path") or "").strip()
    if explicit and os.path.exists(explicit):
        return explicit
    tokens = []
    for name in (target.get("name"), target.get("label"), target.get("key")):
        text = str(name or "").strip()
        if text:
            tokens.extend([text, text.replace(".", "_"), text.replace("_", ".")])
    hits = []
    for token in dict.fromkeys(tokens):
        for pattern in (
            os.path.join(_WBC_DB_FILES, f"*{token}*Device_Deployment.xlsx"),
            os.path.join(_WBC_DB_FILES, f"*{token}*Device_Deployment.xlsm"),
            os.path.join(_WBC_DB_FILES, f"*{token}*.xlsx"),
            os.path.join(_WBC_DB_FILES, f"*{token}*.xlsm"),
        ):
            hits.extend(glob(pattern))
    # Mostly-matched fallback: try base name (first alphabetic characters, min 4 chars)
    # e.g. "Kobuk11" -> "Kobuk", "Pinnacles.2.3" -> "Pinnacles"
    if not hits:
        base_names: set = set()
        for name in (target.get("name"), target.get("label"), target.get("key")):
            text = str(name or "").strip()
            if text:
                m = re.match(r"^([A-Za-z]+)", text)
                if m and len(m.group(1)) >= 4:
                    base_names.add(m.group(1))
        for base in base_names:
            for pattern in (
                os.path.join(_WBC_DB_FILES, f"*{base}*Device_Deployment.xlsx"),
                os.path.join(_WBC_DB_FILES, f"*{base}*Device_Deployment.xlsm"),
                os.path.join(_WBC_DB_FILES, f"*{base}*.xlsx"),
                os.path.join(_WBC_DB_FILES, f"*{base}*.xlsm"),
            ):
                hits.extend(glob(pattern))
    hits = [p for p in dict.fromkeys(hits) if os.path.exists(p)]
    if hits:
        hits.sort(key=lambda p: ("device_deployment" not in os.path.basename(p).lower(), os.path.basename(p).lower()))
        return hits[0]
    return explicit


def _find_target_fr_workbook(target: Dict[str, str], db_cfg: Dict[str, str] = None) -> str:
    """Resolve the PL-wise FR workbook; MTBF remains sourced from _WBC_DB_FILES."""
    explicit = str((db_cfg or {}).get("fr_excel_path") or "").strip()
    if explicit and os.path.exists(explicit):
        return explicit

    # Build a rich set of search tokens from name/label/key
    tokens: List[str] = []
    for value in (target.get("name"), target.get("label"), target.get("key")):
        text = str(value or "").strip()
        if not text:
            continue
        tokens.append(text)
        tokens.append(text.replace(".", "_"))
        tokens.append(text.replace("_", "."))
        # Base name without LE version (e.g. "Tarang" from "Tarang.LE.1.0")
        base = re.split(r"[._]LE[._]", text, maxsplit=1)[0].strip()
        if base and base != text:
            tokens.append(base)
        # Version-compressed token: base + version digits (e.g. "Tarang1.0" from "Tarang.LE.1.0")
        # Handles filenames like Tarang1.0_Device_Deployment.xlsx
        m_ver = re.search(r"[._]LE[._](\d+[._]\d+)", text, re.I)
        if m_ver and base:
            ver = m_ver.group(1).replace("_", ".")
            tokens.append(f"{base}{ver}")          # Tarang1.0
            tokens.append(f"{base}{ver.replace('.', '_')}")  # Tarang1_0
            tokens.append(f"{base}{ver.replace('.', '')}")   # Tarang10

    search_dirs = [d for d in [_WBC_FR_FILES, _WBC_DB_FILES] if d and os.path.isdir(d)]

    hits: List[str] = []
    for search_dir in search_dirs:
        for token in dict.fromkeys(tokens):
            # Strict: require "Device_Deployment" in filename
            hits.extend(glob(os.path.join(search_dir, f"*{token}*Device_Deployment.xls*")))
        # Broader fallback: any workbook containing the token (no suffix requirement)
        if not hits:
            for token in dict.fromkeys(tokens):
                hits.extend(glob(os.path.join(search_dir, f"*{token}*.xls*")))

    hits = [path for path in dict.fromkeys(hits) if os.path.exists(path)]
    if hits:
        hits.sort(key=lambda path: os.path.getmtime(path), reverse=True)
        return hits[0]
    return ""


def _wbc_header_index(headers: List[str], names: List[str]) -> int:
    normalized = [_norm(h) for h in (headers or [])]
    for name in names:
        n = _norm(name)
        if n in normalized:
            return normalized.index(n)
    for i, h in enumerate(normalized):
        if any(_norm(name) and _norm(name) in h for name in names):
            return i
    return -1


def _coerce_wbc_mtbf_payload(data: Dict[str, Any]) -> Dict[str, Any]:
    data = data if isinstance(data, dict) else {}
    headers = data.get("headers") or []
    rows = data.get("rows") or []
    chart_rows = []
    if rows and headers:
        idx_sno = _wbc_header_index(headers, ["S.No", "S No", "sno"])
        idx_build = _wbc_header_index(headers, ["CRM Build ID", "CRM Build", "build id", "build"])
        idx_date = _wbc_header_index(headers, ["Date"])
        idx_hours = _wbc_header_index(headers, ["Hours+", "Hours", "hour"])
        idx_crash = _wbc_header_index(headers, ["crash", "crashes"])
        idx_mtbf = _wbc_header_index(headers, ["MTBF"])
        for i, row in enumerate(rows, start=1):
            vals = row.get("values") or [] if isinstance(row, dict) else []
            def at(idx: int) -> str:
                return str(vals[idx] if 0 <= idx < len(vals) else "").strip()
            build = at(idx_build)
            if not build:
                continue
            crash = _safe_int(at(idx_crash)) if idx_crash >= 0 else 0
            hours = _safe_float(at(idx_hours)) if idx_hours >= 0 else 0.0
            mtbf = _safe_float(at(idx_mtbf)) if idx_mtbf >= 0 else 0.0
            if not mtbf and hours and crash:
                mtbf = round(hours / crash, 2)
            chart_rows.append({
                "id": f"wbc_{i}",
                "s_no": _safe_int(at(idx_sno)) or i,
                "crm_build_id": build,
                "meta_id": build,
                "date": at(idx_date)[:10] if idx_date >= 0 else "",
                "hours": round(hours, 2),
                "crash": crash,
                "total_crashes": crash,
                "mtbf": round(mtbf, 2),
            })
    elif data.get("chart_rows"):
        for i, row in enumerate(data.get("chart_rows") or [], start=1):
            build = str(row.get("crm_build_id") or row.get("meta_id") or row.get("build") or "").strip()
            if not build:
                continue
            crash = _safe_int(row.get("crash") if row.get("crash") not in (None, "") else row.get("total_crashes"))
            chart_rows.append({
                "id": str(row.get("id") or f"wbc_{i}"),
                "s_no": _safe_int(row.get("s_no")) or i,
                "crm_build_id": build,
                "meta_id": build,
                "date": str(row.get("date") or "").strip()[:10],
                "hours": round(_safe_float(row.get("hours")), 2),
                "crash": crash,
                "total_crashes": crash,
                "mtbf": round(_safe_float(row.get("mtbf")), 2),
            })
    data["headers"] = _mtbf_headers()
    data["mtbf_headers"] = _mtbf_headers()
    data["chart_rows"] = chart_rows
    data["rows"] = [{"excel_row": i + 2, "values": [r.get("s_no"), r.get("crm_build_id"), r.get("date"), r.get("hours"), r.get("crash"), r.get("mtbf")], "row": dict(zip(_mtbf_headers(), [r.get("s_no"), r.get("crm_build_id"), r.get("date"), r.get("hours"), r.get("crash"), r.get("mtbf")]))} for i, r in enumerate(chart_rows)]
    data["wbc_mtbf_format"] = True
    return data


def _load_or_sync_mainline_mtbf(target: Dict[str, str], db_cfg: Dict[str, str]) -> Dict[str, Any]:
    key = target.get("key") or target.get("name") or "target"
    mtbf_path = _mtbf_json_path(key)
    data = _read_json(mtbf_path, {})
    if data.get("chart_rows"):
        coerced = _coerce_wbc_mtbf_payload(data)
        # Sync to internal dashboard JSON on first access (or whenever it is missing).
        # This ensures the internal MTBF page shows the same data as the WBC Live View
        # Stats page without requiring a manual save/sync operation.
        internal_path = _dashboard_mtbf_json_path(key, target)
        if not os.path.exists(internal_path):
            _sync_to_dashboard_mtbf_json(key, coerced, target=target)
        return coerced
    excel_path = _find_target_excel(target, db_cfg)
    if not excel_path or not os.path.exists(excel_path):
        return {"headers": [], "rows": [], "chart_rows": [], "error": f"Mainline_Build_Details Excel not found for {target.get('label') or key}", "excel_path": excel_path}
    try:
        sheets = _workbook_sheets(excel_path)
        sheet = next((s for s in sheets if str(s).strip().lower() == "mainline_build_details"), "")
        if not sheet:
            sheet = next((s for s in sheets if "mainline" in str(s).lower() and "build" in str(s).lower()), "")
        if not sheet:
            raise ValueError("Sheet Mainline_Build_Details not found")
        data = _coerce_wbc_mtbf_payload(_sheet_to_payload(key, excel_path, sheet))
        data["target_key"] = key
        data["target_label"] = target.get("label") or target.get("name") or key
        data["mtbf_sheet"] = sheet
        data["one_time_synced"] = True
        data["saved_json"] = mtbf_path
        _write_json(mtbf_path, data)
        _write_json(_target_json_path(key), data)
        _sync_to_dashboard_mtbf_json(key, data, target=target)
        return data
    except Exception as exc:
        return {"headers": [], "rows": [], "chart_rows": [], "error": str(exc), "excel_path": excel_path, "sheet_name": "Mainline_Build_Details"}


def _discover_targets(include_hidden: bool = False) -> List[Dict[str, str]]:
    files = []
    for pattern in (os.path.join(_WBC_DB_FILES, "*_Device_Deployment.xlsx"), os.path.join(_WBC_DB_FILES, "*_Device_Deployment.xlsm")):
        files.extend(glob(pattern))
    files = [p for p in files if not os.path.basename(p).startswith("~$")]
    rows = [_target_from_excel(p) for p in sorted(dict.fromkeys(files), key=lambda x: os.path.basename(x).lower())]

    # Also expose manually configured WBC targets. This allows WBC to behave like
    # HGY/HQX dashboards where a target can be added directly by selecting the
    # required JIRAs/Open JIRAs/CR tables, even when no deployment workbook exists.
    cfg = _read_json(_config_path(), {})
    hidden = {_slug(x).lower() for x in (cfg.get("hidden_targets") or [])}
    if not include_hidden:
        rows = [r for r in rows if str(r.get("key") or "").lower() not in hidden]
    by_key = {str(r.get("key") or "").lower(): r for r in rows}
    for key, row in (cfg.get("targets") or {}).items():
        row = row if isinstance(row, dict) else {}
        target_key = _slug(key)
        low = target_key.lower()
        if low in hidden and not include_hidden:
            continue
        if low in by_key:
            by_key[low].update({
                "name": row.get("name") or by_key[low].get("name") or target_key,
                "label": _friendly_label(row.get("label") or by_key[low].get("label") or target_key),
                "manual": bool(row.get("manual")),
            })
            continue
        name = str(row.get("name") or row.get("label") or target_key).strip() or target_key
        by_key[low] = {
            "key": target_key,
            "name": name,
            "label": _friendly_label(str(row.get("label") or name).strip() or name),
            "excel_path": str(row.get("excel_path") or "").strip(),
            "manual": True,
        }
    out = sorted(by_key.values(), key=lambda r: str(r.get("label") or r.get("key") or "").lower())
    if include_hidden:
        for row in out:
            row["hidden"] = str(row.get("key") or "").lower() in hidden
    return out


def _load_config() -> Dict[str, Any]:
    cfg = _read_json(_config_path(), {})
    cfg.setdefault("root", _WBC_ROOT)
    cfg.setdefault("db_files", _WBC_DB_FILES)
    cfg.setdefault("targets", {})
    cfg.setdefault("hidden_targets", [])
    cfg.setdefault("updated_at", "")
    return cfg


def _save_config(payload: Dict[str, Any]) -> Dict[str, Any]:
    cfg = _load_config()
    if isinstance(payload.get("hidden_targets"), list):
        cfg["hidden_targets"] = [_slug(x) for x in payload.get("hidden_targets") if str(x or "").strip()]
    if isinstance(payload.get("targets"), dict):
        cleaned = {}
        for key, row in payload.get("targets", {}).items():
            if not str(key or "").strip():
                continue
            row = row if isinstance(row, dict) else {}
            target_key = _slug(str(key))
            cleaned[target_key] = {
                "name": str(row.get("name") or row.get("label") or target_key).strip(),
                "label": _friendly_label(str(row.get("label") or row.get("name") or target_key).strip()),
                "excel_path": str(row.get("excel_path") or "").strip(),
                "manual": bool(row.get("manual") or not str(row.get("excel_path") or "").strip()),
                "jiras_table": str(row.get("jiras_table") or row.get("target_table") or "").strip(),
                "target_table": str(row.get("target_table") or row.get("jiras_table") or "").strip(),
                "openjiras_table": str(row.get("openjiras_table") or "").strip(),
                "unique_crs_table": str(row.get("unique_crs_table") or "").strip(),
                "overall_crs_table": str(row.get("overall_crs_table") or "").strip(),
            }
        cfg["targets"] = cleaned
    cfg["updated_at"] = datetime.utcnow().isoformat() + "Z"
    _write_json(_config_path(), cfg)
    return cfg


def _bt(schema: str, table: str) -> str:
    return f"`{str(schema).strip('`')}`.`{str(table).strip('`')}`"


def _split_table(value: str) -> Tuple[str, str]:
    text = str(value or "").replace("`", "").strip()
    if "." in text:
        schema, table = text.split(".", 1)
        return schema.strip(), table.strip()
    return _WBC_SCHEMA, text


def _table_cols(cur, fq_table: str) -> List[str]:
    schema, table = _split_table(fq_table)
    if not table:
        return []
    cur.execute("SELECT 1 FROM information_schema.tables WHERE table_schema=%s AND table_name=%s LIMIT 1", (schema, table))
    if cur.fetchone() is None:
        return []
    cur.execute(f"SHOW COLUMNS FROM {_bt(schema, table)}")
    return [str(r.get("Field") or "") for r in (cur.fetchall() or []) if r.get("Field")]


def _norm(value: Any) -> str:
    return re.sub(r"[^a-z0-9]+", "", str(value or "").lower())


def _first_col(cols: List[str], candidates: List[str]) -> str:
    by_norm = {_norm(c): c for c in cols}
    for cand in candidates:
        if by_norm.get(_norm(cand)):
            return by_norm[_norm(cand)]
    for col in cols:
        if any(_norm(cand) and _norm(cand) in _norm(col) for cand in candidates):
            return col
    return ""


def _safe_int(value: Any) -> int:
    try:
        return int(float(str(value or 0).replace(",", "").strip()))
    except Exception:
        return 0


def _safe_float(value: Any) -> float:
    try:
        return float(str(value or 0).replace(",", "").strip())
    except Exception:
        return 0.0


def _db_table_options() -> List[Dict[str, str]]:
    conn = get_mysql_connection_db(database_name=_WBC_SCHEMA) or get_mysql_connection_db(bu_key=None)
    if not conn:
        return []
    cur = conn.cursor(dictionary=True)
    try:
        cur.execute(
            """
            SELECT TABLE_NAME FROM information_schema.TABLES
            WHERE TABLE_SCHEMA=%s
              AND (TABLE_NAME LIKE '%%_jiras' OR TABLE_NAME LIKE '%%_openjiras'
                   OR TABLE_NAME LIKE '%%_unique_crs' OR TABLE_NAME LIKE '%%overall%%cr%%')
            ORDER BY CASE
              WHEN TABLE_NAME LIKE '%%_jiras' THEN 0
              WHEN TABLE_NAME LIKE '%%_openjiras' THEN 1
              WHEN TABLE_NAME LIKE '%%_unique_crs' THEN 2
              WHEN TABLE_NAME LIKE '%%overall%%cr%%' THEN 3 ELSE 9 END, TABLE_NAME
            LIMIT 2000
            """,
            (_WBC_SCHEMA,),
        )
        out = []
        for row in cur.fetchall() or []:
            name = str(row.get("TABLE_NAME") or "").strip()
            low = name.lower()
            kind = "other"
            if low.endswith("_openjiras"):
                kind = "openjiras"
            elif low.endswith("_unique_crs"):
                kind = "unique_crs"
            elif "overall" in low and "cr" in low:
                kind = "overallcrs"
            elif low.endswith("_jiras"):
                kind = "jiras"
            out.append({"name": name, "fq": f"`{_WBC_SCHEMA}`.`{name}`", "kind": kind, "label": f"{name} ({kind})"})
        return out
    except Exception:
        return []
    finally:
        try:
            cur.close(); conn.close()
        except Exception:
            pass


def _count_from_table(fq_table: str, candidates: List[str]) -> int:
    if not fq_table:
        return 0
    conn = get_mysql_connection_db(database_name=_WBC_SCHEMA) or get_mysql_connection_db(bu_key=None)
    if not conn:
        return 0
    cur = conn.cursor(dictionary=True)
    try:
        cols = _table_cols(cur, fq_table)
        if not cols:
            return 0
        col = _first_col(cols, candidates) or cols[0]
        schema, table = _split_table(fq_table)
        cur.execute(f"SELECT COUNT(DISTINCT NULLIF(TRIM(`{col}`), '')) AS cnt FROM {_bt(schema, table)}")
        return _safe_int((cur.fetchone() or {}).get("cnt"))
    except Exception:
        return 0
    finally:
        try:
            cur.close(); conn.close()
        except Exception:
            pass


def _target_pl_candidates(target: Dict[str, str]) -> List[str]:
    vals: List[str] = []
    for raw in (target.get("label"), target.get("name"), target.get("key")):
        text = str(raw or "").strip()
        if not text:
            continue
        vals.append(text)
        compact = re.sub(r"[^A-Za-z0-9]", "", text)
        m = re.match(r"(?i)^([A-Za-z]+)(\d)(\d)$", compact)
        if m:
            vals.append(f"{m.group(1)}.LE.{m.group(2)}.{m.group(3)}")
        m2 = re.match(r"(?i)^([A-Za-z]+)LE(\d)(\d)$", compact)
        if m2:
            vals.append(f"{m2.group(1)}.LE.{m2.group(2)}.{m2.group(3)}")
    seen, out = set(), []
    for val in vals:
        key = val.upper()
        if key not in seen:
            seen.add(key); out.append(val)
    return out


def _pl_values(cur, fq_table: str, target: Dict[str, str] = None) -> List[str]:
    cols = _table_cols(cur, fq_table)
    if not cols:
        return []
    # Use only real PL/product columns. Do not use broad target/program columns because
    # those can match every Kobuk/Kuna running meta and inflate WBC current builds.
    col = _first_col(cols, ["pl_id", "PL-ID", "PL ID", "product_line", "software_product", "cpl"])
    if not col:
        return []
    target_candidates = _target_pl_candidates(target or {})
    target_keys = {_norm(str(x)) for x in target_candidates if str(x or "").strip()}
    schema, table = _split_table(fq_table)
    cur.execute(f"SELECT DISTINCT `{col}` AS v FROM {_bt(schema, table)} WHERE `{col}` IS NOT NULL AND TRIM(`{col}`)<>'' LIMIT 300")
    seen, out = set(), []
    for row in cur.fetchall() or []:
        val = str(row.get("v") or "").strip()
        if not val:
            continue
        if target_keys and _norm(val) not in target_keys:
            continue
        key = val.upper()
        if key not in seen:
            seen.add(key); out.append(val)
    return out or target_candidates


def _meta_label(build: str) -> str:
    text = str(build or "").split("/")[-1].split("\\")[-1]
    m = re.search(r"(?i)meta[-_ ]?0*(\d{2,6})", text) or re.search(r"-0*(\d{3,6})(?:[.-]|$)", text)
    return f"Meta-{int(m.group(1)):04d}" if m else (text[:50] or "-")


def _current_running_builds(target: Dict[str, str], db_cfg: Dict[str, str]) -> Dict[str, Any]:
    conn = get_mysql_connection_db(bu_key=None)
    if not conn:
        return {"rows": [], "updated_at": "", "source": "No DB connection"}
    cur = conn.cursor(dictionary=True)
    try:
        terms = []
        for fq in (db_cfg.get("jiras_table"), db_cfg.get("openjiras_table"), db_cfg.get("unique_crs_table")):
            if fq:
                terms.extend(_pl_values(cur, fq, target))
        if not terms:
            terms.extend(_target_pl_candidates(target))
        seen, wheres, params = set(), [], []
        for term in terms:
            term = str(term or "").strip()
            key = term.upper()
            if not term or key in seen:
                continue
            seen.add(key)
            wheres.append("(software_product = %s OR product_flavor = %s OR build_name LIKE %s OR build_id LIKE %s)")
            params.extend([term, term, f"%{term}%", f"%{term}%"])
        where_sql = " OR ".join(wheres) if wheres else "1=0"
        cur.execute("SELECT MAX(updated_at) AS updated_at FROM pdt_stats_dashboard.axiom_job_summary")
        meta = cur.fetchone() or {}
        cur.execute(
            f"""
            SELECT build_id, build_name, software_product, product_flavor, device_count, job_id, started_at, submitted_at, hours
            FROM pdt_stats_dashboard.axiom_job_summary
            WHERE state='Running' AND ({where_sql})
            ORDER BY submitted_at DESC LIMIT 300
            """,
            tuple(params),
        )
        grouped: Dict[str, Dict[str, Any]] = {}
        for row in cur.fetchall() or []:
            build = str(row.get("build_name") or row.get("build_id") or "").strip().split("/")[-1].split("\\")[-1]
            if not build:
                continue
            item = grouped.setdefault(build, {
                "build_id": build, "meta_id": _meta_label(build), "job_count": 0, "device_count": 0,
                "hours": 0.0, "software_product": str(row.get("software_product") or ""),
                "product_flavor": str(row.get("product_flavor") or ""), "started_at": str(row.get("started_at") or "")[:19],
            })
            item["job_count"] += 1
            item["device_count"] = max(_safe_int(item.get("device_count")), _safe_int(row.get("device_count")))
            item["hours"] = round(_safe_float(item.get("hours")) + _safe_float(row.get("hours")), 2)
        rows = list(grouped.values())
        rows.sort(key=lambda r: str(r.get("started_at") or ""), reverse=True)
        return {"rows": rows, "updated_at": str(meta.get("updated_at") or ""), "source": "axiom_job_summary"}
    except Exception as exc:
        return {"rows": [], "updated_at": "", "source": str(exc)}
    finally:
        try:
            cur.close(); conn.close()
        except Exception:
            pass


def _where_open_cr(cols: List[str]) -> str:
    status_col = _first_col(cols, ["cr_status", "status", "state"])
    if not status_col:
        return ""
    return f" WHERE LOWER(`{status_col}`) REGEXP 'open|analysis|new|assigned|inprogress|in_progress|fix'"


def _preview_rows_filtered(fq_table: str, limit: int = 100, open_cr_only: bool = False) -> Dict[str, Any]:
    if not fq_table:
        return {"columns": [], "rows": [], "count": 0, "error": "No table configured"}
    conn = get_mysql_connection_db(database_name=_WBC_SCHEMA) or get_mysql_connection_db(bu_key=None)
    if not conn:
        return {"columns": [], "rows": [], "count": 0, "error": "No DB connection"}
    cur = conn.cursor(dictionary=True)
    try:
        cols = _table_cols(cur, fq_table)
        if not cols:
            return {"columns": [], "rows": [], "count": 0, "error": "Table not found"}
        schema, table = _split_table(fq_table)
        selected = cols[:18]
        where_sql = _where_open_cr(cols) if open_cr_only else ""
        cur.execute(f"SELECT COUNT(*) AS cnt FROM {_bt(schema, table)}{where_sql}")
        total = _safe_int((cur.fetchone() or {}).get("cnt"))
        cur.execute(f"SELECT {', '.join('`'+c+'`' for c in selected)} FROM {_bt(schema, table)}{where_sql} LIMIT %s", (limit,))
        rows = [{k: (v.isoformat() if isinstance(v, (date, datetime)) else ("" if v is None else v)) for k, v in r.items()} for r in (cur.fetchall() or [])]
        return {"columns": selected, "rows": rows, "count": total, "error": ""}
    except Exception as exc:
        return {"columns": [], "rows": [], "count": 0, "error": str(exc)}
    finally:
        try:
            cur.close(); conn.close()
        except Exception:
            pass


def _build_tail(value: Any) -> str:
    return str(value or "").strip().replace("/", "\\").split("\\")[-1]


def _norm_build_key(value: Any) -> str:
    text = _build_tail(value).upper()
    return re.sub(r"[^A-Z0-9]+", "", text)


def _build_match_tokens(value: Any) -> set:
    raw = _build_tail(value).upper()
    norm = _norm_build_key(raw)
    tokens = {norm} if norm else set()
    for m in re.findall(r"\b\d{3,6}\b", raw):
        tokens.add(m.lstrip("0") or m)
    for m in re.findall(r"(?i)(KOBUK|KUNA|PINEAPPLE|QCM|QCS|TURING|TARANG)[A-Z0-9_.-]*", raw):
        tokens.add(_norm_build_key(m))
    for m in re.findall(r"(?i)(ROOT[A-Z0-9_.-]+|NON[A-Z0-9_.-]+|TAURIN[A-Z0-9_.-]+|CPFWK[A-Z0-9_.-]+)", raw):
        tokens.add(_norm_build_key(m))
    return {t for t in tokens if t}


def _builds_match(a: Any, b: Any) -> bool:
    ak, bk = _norm_build_key(a), _norm_build_key(b)
    if not ak or not bk:
        return False
    if ak == bk or ak in bk or bk in ak:
        return True
    at, bt = _build_match_tokens(a), _build_match_tokens(b)
    common = at.intersection(bt)
    if not common:
        return False
    # A meta/build number alone is too weak. Require at least one additional
    # product/flavor token or a strong containment match above.
    strong = [t for t in common if not re.fullmatch(r"\d{1,6}", t)]
    return bool(strong) or len(common) >= 2


def _build_summary_from_jiras(jiras_table: str, openjiras_table: str = "") -> Dict[str, Any]:
    """Build-wise consolidated report for WBC, matching the HGY/HQX tab flow.

    This report is intentionally table-based and covers all historical builds in
    the configured JIRAs/Open JIRAs tables. It is separate from Current Running
    Builds, which uses Axiom/saved JQL cards.
    """
    source_defs = [
        ("JIRAs", str(jiras_table or "").strip()),
        ("Open JIRAs", str(openjiras_table or "").strip()),
    ]
    sources = [(label, table) for label, table in source_defs if table]
    if not sources:
        return {"builds": [], "rows_by_build": {}, "error": "No JIRAs/Open JIRAs table configured"}
    conn = get_mysql_connection_db(database_name=_WBC_SCHEMA) or get_mysql_connection_db(bu_key=None)
    if not conn:
        return {"builds": [], "rows_by_build": {}, "error": "No DB connection"}
    cur = conn.cursor(dictionary=True)
    try:
        grouped: Dict[str, Dict[str, Any]] = {}
        rows_by_build: Dict[str, List[Dict[str, Any]]] = {}
        errors: List[str] = []
        for source_label, source in sources:
            cols = _table_cols(cur, source)
            if not cols:
                errors.append(f"Table not found: {source}")
                continue
            build_col = _first_col(cols, ["metabuild", "MetaBuild", "meta_build", "build_id", "build_name", "build", "builds", "CRM Build ID", "Meta-ID"])
            jira_col = _first_col(cols, ["stability_ticket", "jira", "jira_id", "jira_key", "ticket", "key"])
            cr_col = _first_col(cols, ["mapped_cr", "mapped_crs", "cr", "cr_id", "crid", "cr_number", "cr_current_ticket", "Change Request"])
            title_col = _first_col(cols, ["jira_title", "title", "summary", "cr_title"])
            area_col = _first_col(cols, ["cr_area", "area", "technology_area", "ChangeRequestParticipant.Area"])
            status_col = _first_col(cols, ["cr_status", "jira_status", "status", "state"])
            date_col = _first_col(cols, ["jira_date", "last_instance", "si_last_seen", "updated", "created", "created_date"])
            if not build_col:
                errors.append(f"No build/metabuild column in {source}")
                continue
            selected = [c for c in [build_col, jira_col, cr_col, title_col, area_col, status_col, date_col] if c]
            for col in cols[:24]:
                if col not in selected:
                    selected.append(col)
            schema, table = _split_table(source)
            order_sql = f" ORDER BY `{date_col}` DESC" if date_col else ""
            cur.execute(
                f"SELECT {', '.join('`'+c+'`' for c in selected[:26])} FROM {_bt(schema, table)} "
                f"WHERE `{build_col}` IS NOT NULL AND TRIM(`{build_col}`)<>''{order_sql} LIMIT 50000"
            )
            for row in cur.fetchall() or []:
                row = {k: (v.isoformat() if isinstance(v, (date, datetime)) else ("" if v is None else v)) for k, v in row.items()}
                build = str(row.get(build_col) or "").strip() or "Unknown Build"
                item = grouped.setdefault(build, {
                    "build_id": build,
                    "meta_id": _meta_label(build),
                    "row_count": 0,
                    "open_jira_count": 0,
                    "area": "",
                    "last_seen": "",
                    "_crs": set(),
                    "_jiras": set(),
                })
                clean_row = {k: v for k, v in row.items() if k not in ("crash_type", "crash_types", "type", "failure_type", "_source_table")}
                clean_row.setdefault("Source", source_label)
                rows_by_build.setdefault(build, []).append(clean_row)
                item["row_count"] += 1
                if source == openjiras_table:
                    item["open_jira_count"] += 1
                cr = str(row.get(cr_col) if cr_col else "").strip()
                jira = str(row.get(jira_col) if jira_col else "").strip()
                if cr:
                    item["_crs"].add(cr)
                if jira:
                    item["_jiras"].add(jira)
                area = str(row.get(area_col) if area_col else "").strip()
                if area and not item.get("area"):
                    item["area"] = area
                seen = str(row.get(date_col) if date_col else "").strip()
                if seen and seen > str(item.get("last_seen") or ""):
                    item["last_seen"] = seen
        builds = []
        for item in grouped.values():
            item["cr_count"] = len(item.pop("_crs", set()))
            item["jira_count"] = len(item.pop("_jiras", set()))
            item["total_count"] = item.get("row_count", 0)
            builds.append(item)
        builds.sort(key=lambda r: (str(r.get("last_seen") or ""), _safe_int(r.get("row_count"))), reverse=True)
        return {
            "ok": True,
            "source": "Configured WBC JIRAs/Open JIRAs tables",
            "builds": builds[:500],
            "rows_by_build": rows_by_build,
            "error": "; ".join(errors),
        }
    except Exception as exc:
        return {"builds": [], "rows_by_build": {}, "error": str(exc)}
    finally:
        try:
            cur.close(); conn.close()
        except Exception:
            pass



def _preview_rows(fq_table: str, limit: int = 100) -> Dict[str, Any]:
    if not fq_table:
        return {"columns": [], "rows": [], "count": 0, "error": "No table configured"}
    conn = get_mysql_connection_db(database_name=_WBC_SCHEMA) or get_mysql_connection_db(bu_key=None)
    if not conn:
        return {"columns": [], "rows": [], "count": 0, "error": "No DB connection"}
    cur = conn.cursor(dictionary=True)
    try:
        cols = _table_cols(cur, fq_table)
        if not cols:
            return {"columns": [], "rows": [], "count": 0, "error": "Table not found"}
        schema, table = _split_table(fq_table)
        selected = cols[:16]
        cur.execute(f"SELECT {', '.join('`'+c+'`' for c in selected)} FROM {_bt(schema, table)} LIMIT %s", (limit,))
        rows = [{k: (v.isoformat() if isinstance(v, (date, datetime)) else ("" if v is None else v)) for k, v in r.items()} for r in (cur.fetchall() or [])]
        return {"columns": selected, "rows": rows, "count": len(rows), "error": ""}
    except Exception as exc:
        return {"columns": [], "rows": [], "count": 0, "error": str(exc)}
    finally:
        try:
            cur.close(); conn.close()
        except Exception:
            pass


def _wbc_pdt_key(target: Dict[str, str]) -> str:
    """
    Return the PDT TARGETS_CONFIG key for a WBC target so that
    live_view_saved_jql_service can resolve the correct BU and write
    cache files to the right path instead of UNKNOWN_BU.

    Strategy (first match wins):
      1. target['pdt_key']  - explicitly set in wbc_config.json
      2. Fuzzy match: find a TARGETS_CONFIG key whose slug matches the
         WBC target key (e.g. 'Kobuk11' -> 'Kobuk.LE.1.1')
      3. Fall back to the WBC key itself (may still produce UNKNOWN_BU
         if not in TARGETS_CONFIG, but that is the existing behaviour).
    """
    wbc_key = str(target.get("key") or "").strip()
    # 1. Explicit override
    if target.get("pdt_key"):
        return str(target["pdt_key"]).strip()
    # 2. Fuzzy: normalise both sides and compare
    try:
        from dashboard_common import get_targets_config
        targets_cfg = get_targets_config() or {}
        wbc_slug = _slug(wbc_key).lower()
        for pdt_key in targets_cfg:
            if _slug(pdt_key).lower() == wbc_slug:
                return pdt_key
        # Also try label match (e.g. 'Kobuk.LE.1.1' label contains 'Kobuk')
        label = str(target.get("label") or target.get("name") or "").lower()
        for pdt_key in targets_cfg:
            if _slug(pdt_key).lower() in label or label.startswith(_slug(pdt_key).lower()):
                return pdt_key
    except Exception:
        pass
    # 3. Fall back
    return wbc_key


def _find_target(target_key: str) -> Dict[str, str]:
    return next((t for t in _discover_targets() if t["key"].lower() == str(target_key or "").lower()), {})


def _mtbf_headers() -> List[str]:
    return ["S.No", "CRM Build ID", "Date", "Hours+", "crash", "MTBF"]


def _normalize_mtbf_row(row: Dict[str, Any], index: int) -> Dict[str, Any]:
    row = row if isinstance(row, dict) else {}
    meta = str(row.get("crm_build_id") or row.get("meta_id") or row.get("build_s") or row.get("build") or row.get("CRM Build ID") or row.get("Meta-ID") or row.get("Build") or "").strip()
    hours = _safe_float(row.get("hours") or row.get("Hours+") or row.get("Hours"))
    total = _safe_int(row.get("crash") if row.get("crash") not in (None, "") else row.get("total_crashes") if row.get("total_crashes") not in (None, "") else row.get("crash") or row.get("Total Crashes"))
    mtbf_raw = row.get("mtbf") if row.get("mtbf") not in (None, "") else row.get("MTBF")
    mtbf = _safe_float(mtbf_raw)
    if not mtbf and hours and total:
        mtbf = round(hours / total, 2)
    return {
        "id": str(row.get("id") or f"manual_{index}"),
        "s_no": index,
        "date": str(row.get("date") or row.get("Date") or "").strip()[:10],
        "crm_build_id": meta,
        "meta_id": meta,
        "hours": round(hours, 2),
        "crash": total,
        "total_crashes": total,
        "mtbf": round(mtbf, 2),
    }


def _save_mtbf_chart_rows(target: Dict[str, str], db_cfg: Dict[str, str], rows: List[Dict[str, Any]]) -> Dict[str, Any]:
    key = target.get("key") or "target"
    data = _load_or_sync_mainline_mtbf(target, db_cfg)
    chart_rows = [_normalize_mtbf_row(r, i + 1) for i, r in enumerate(rows or []) if isinstance(r, dict)]
    headers = _mtbf_headers()
    data.update({
        "target": key,
        "target_key": key,
        "target_label": target.get("label") or target.get("name") or key,
        "sheet_name": data.get("sheet_name") or "Mainline_Build_Details",
        "mtbf_sheet": data.get("mtbf_sheet") or "Mainline_Build_Details",
        "headers": headers,
        "mtbf_headers": headers,
        "chart_rows": chart_rows,
        "rows": [{"excel_row": i + 2, "values": [r.get("s_no"), r.get("crm_build_id"), r.get("date"), r.get("hours"), r.get("crash"), r.get("mtbf")], "row": dict(zip(headers, [r.get("s_no"), r.get("crm_build_id"), r.get("date"), r.get("hours"), r.get("crash"), r.get("mtbf")]))} for i, r in enumerate(chart_rows)],
        "updated_at": datetime.utcnow().isoformat() + "Z",
        "edited_json": True,
        "saved_json": _mtbf_json_path(key),
    })
    _write_json(_mtbf_json_path(key), data)
    _write_json(_target_json_path(key), data)
    _sync_to_dashboard_mtbf_json(key, data, target=target)
    return data


def _read_swpdt_summary_from_excel(excel_path: str) -> Dict[str, str]:
    """Read SWPDT_Summary (→ overview) and SWPDT_Summary2 (→ pdt_status) from the workbook."""
    import openpyxl
    result: Dict[str, str] = {"overview": "", "pdt_status": ""}
    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        # SWPDT_Summary: col 2 rows 2+ = summary bullet points
        if "SWPDT_Summary" in wb.sheetnames:
            ws = wb["SWPDT_Summary"]
            lines = []
            for r in range(2, ws.max_row + 1):
                val = str(ws.cell(r, 2).value or "").strip().replace("\xa0", " ").strip()
                if val:
                    lines.append(val)
            result["overview"] = "\n".join(lines)
        # SWPDT_Summary2: col 1 rows 2+ = PDT status / timelines
        if "SWPDT_Summary2" in wb.sheetnames:
            ws2 = wb["SWPDT_Summary2"]
            lines2 = []
            for r in range(2, ws2.max_row + 1):
                val = str(ws2.cell(r, 1).value or "").strip().replace("\xa0", " ").strip()
                if val:
                    lines2.append(val)
            result["pdt_status"] = "\n".join(lines2)
        wb.close()
    except Exception:
        pass
    return result


def _load_overview_summary(target_key: str) -> Dict[str, Any]:
    data = _read_json(_overview_summary_path(target_key), {})
    if not isinstance(data, dict):
        data = {}
    data.setdefault("target_key", _slug(target_key))
    data.setdefault("engineer", "")
    data.setdefault("summary_title", "")
    data.setdefault("overview", "")
    data.setdefault("overview_html", "")
    data.setdefault("highlights", "")
    data.setdefault("highlights_html", "")
    data.setdefault("risks", "")
    data.setdefault("pdt_status", data.get("next_steps") or "")
    data.setdefault("pdt_status_html", "")
    data.setdefault("next_steps", "")
    data.setdefault("next_steps_html", "")
    data.setdefault("updated_at", "")
    data.setdefault("updated_by", "")
    return data


def _save_overview_summary(target_key: str, payload: Dict[str, Any]) -> Dict[str, Any]:
    payload = payload if isinstance(payload, dict) else {}
    allowed = ["engineer", "summary_title", "overview", "overview_html", "highlights", "highlights_html", "risks", "pdt_status", "pdt_status_html", "next_steps", "next_steps_html"]
    data = _load_overview_summary(target_key)
    for key in allowed:
        data[key] = str(payload.get(key) or "").strip()
    data["target_key"] = _slug(target_key)
    data["updated_at"] = datetime.utcnow().isoformat() + "Z"
    data["updated_by"] = str(getattr(current_user, "id", "") or getattr(current_user, "username", "") or "").strip()
    _write_json(_overview_summary_path(target_key), data)
    return data


def _jql_quote(value: Any) -> str:
    return '"' + str(value or "").replace("\\", "\\\\").replace('"', '\\"') + '"'


def _wbc_build_jql(build_id: str) -> str:
    build = _build_tail(build_id)
    return (
        f'(summary ~ {_jql_quote(build)}) '
        f'AND filter = {JIRA_PDT_FILTER_ID} '
        f'AND (project = QSTABILITY OR project = DROIDBUG OR project = CHIPMD) '
        f'AND summary !~ "tombstone" ORDER BY created ASC'
    )


def _parse_iso_dt(value: Any) -> datetime:
    text = str(value or "").strip().replace("Z", "")
    if not text:
        return datetime.min
    try:
        return datetime.fromisoformat(text)
    except Exception:
        return datetime.min


def _wbc_external_saved_jql_row(row: Dict[str, Any]) -> Dict[str, Any]:
    """Return only current-build metadata safe for external/view-only users."""
    allowed = (
        "id", "name", "build_id", "has_cached_report", "cached_report_stale",
        "last_run_at", "next_run_at", "cache_ttl_minutes", "cached_row_count",
        "cached_cr_count", "cached_jira_count", "cache_status",
    )
    return {key: row.get(key) for key in allowed if key in row}


def _wbc_external_saved_jql_report(report: Dict[str, Any]) -> Dict[str, Any]:
    """Remove saved-filter/JQL implementation details from an external response."""
    hidden_keys = {
        "jql", "raw_jql", "resolved_jql", "filter_id", "filter_resolved",
        "filter_error", "tab",
    }
    return {key: value for key, value in report.items() if key not in hidden_keys}


def _wbc_saved_jql_cache_meta(cached: Dict[str, Any]) -> Dict[str, Any]:
    """Return normalized cache timing/count metadata for a saved-JQL report."""
    cached = cached if isinstance(cached, dict) else {}
    ttl = timedelta(minutes=30)
    generated_at = _parse_iso_dt(cached.get("generated_at"))
    now = datetime.utcnow()
    next_run = generated_at + ttl if generated_at != datetime.min else datetime.min
    expired = bool(generated_at != datetime.min and now >= next_run)
    rows = cached.get("rows") or cached.get("flat_rows") or []
    return {
        "has_cached_report": bool(cached),
        "cached_report_stale": expired,
        "last_run_at": cached.get("generated_at") or "",
        "next_run_at": next_run.isoformat() + "Z" if next_run != datetime.min else "",
        "cache_ttl_minutes": 30,
        "cached_row_count": _safe_int(cached.get("row_count", len(rows))) if cached else 0,
        "cached_cr_count": _safe_int(cached.get("cr_count")) if cached else 0,
        "cached_jira_count": _safe_int(cached.get("jira_count")) if cached else 0,
        "cache_status": "stale" if expired else ("cached" if cached else "not_run"),
    }
    return rows



# ---------------------------------------------------------------------------
# CR / JIRA classification helpers
# Mirrors the CR_EQ list used in auto_gen45_live_view_stats.html and
# live_status_publish_edit_nonau.html so WBC uses the same logic everywhere.
# ---------------------------------------------------------------------------
# Invalid resolution keywords — mirrors _WBC_INV_KW in the frontend
_INVALID_KEYWORDS = [
    "invalid", "incomplete", "rejected", "won't fix", "wont fix",
    "cannot reproduce", "not a bug", "no sir", "nosir", "not applicable",
    "obsolete", "postponed", "withdrawn", "cannotduplicate", "invalid_dup",
    "setup issue", "setup_issue", "incomplete ram dump",
    "similar substring", "substring found",
]


def _is_invalid_row(r: Dict[str, Any]) -> bool:
    """Return True if a flattened row has an invalid resolution/status.
    Checks JIRA Resolution, Final Resolution, Final Status, JIRA Status,
    Resolution Notes and CR Status — same fields as _isInvalidRow() in JS.
    """
    hay = " ".join([
        str(r.get("JIRA Resolution") or ""),
        str(r.get("Final Resolution") or ""),
        str(r.get("Final Status") or ""),
        str(r.get("JIRA Status") or ""),
        str(r.get("Resolution Notes") or ""),
        str(r.get("CR Status") or ""),
    ]).lower()
    return any(kw in hay for kw in _INVALID_KEYWORDS)


_CR_EQUIV_PREFIXES = {
    "ADSPIMAGE", "CNSSDEBUG", "CHIPMD", "ADSPBUG", "CNSS", "WLAN", "QWINBUG",
    "ARAST", "AVATAR", "AVATARWPAP", "BAGHEERAST", "BLAUNCH", "DINOSTABLE", "DROIDBUG",
    "ELANSTABLE", "FORINO", "FRODOST", "FUSIONT", "FUSNFOURST", "JINGALA", "QNPSTBLT",
    "QSTABILITY", "TORINOST", "WAVEAPOLLO", "WCNSTABLE", "WPARAGORN", "WPFRODO",
    "WRSTABLE", "UIBUG", "RMASLT", "SCSTABLE", "AISW", "WPST",
}


def _is_true_cr(value: str) -> bool:
    """Return True if value is a real Orbit CR number (CR followed by 5-9 digits)."""
    return bool(re.match(r"^CR\d{5,9}$", str(value or "").strip(), re.I))


def _is_cr_equiv(value: str) -> bool:
    """Return True if value is a JIRA key from a CR-equivalent project (e.g. WPST-1234)."""
    upper = str(value or "").strip().upper()
    project = upper.split("-")[0] if "-" in upper else ""
    return bool(project and project in _CR_EQUIV_PREFIXES)


def _row_type(cr: str) -> str:
    """Classify a CR/mapped-ticket value into 'cr', 'mapped_jira', or 'open'."""
    if _is_true_cr(cr):
        return "cr"
    if _is_cr_equiv(cr):
        return "mapped_jira"
    return "open"


def _wbc_flatten_consolidated_report(report: Dict[str, Any]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    for cr_row in (report.get("hierarchical_report") or []):
        cr = cr_row.get("cr") or "NO_CR"
        rtype = _row_type(cr)
        jiras = cr_row.get("jiras") or []
        if not jiras:
            rows.append({
                "Row Type": rtype,
                "CR": cr,
                "CR Title": cr_row.get("cr_title") or "",
                "CR Status": cr_row.get("cr_status") or "",
                "CR Area": cr_row.get("cr_area") or "",
                "CR Subsystem": cr_row.get("cr_subsystem") or "",
                "CR Function": cr_row.get("cr_function") or "",
                "JIRA": "",
                "JIRA Title": "",
                "JIRA Status": "",
                "Final Ticket": "",
                "Final Status": "",
                "Final Resolution": "",
                "Mapping Type": "",
                "Resolution Notes": "",
            })
            continue
        for jira in jiras:
            trav = jira.get("traversal") or {}
            rows.append({
                "Row Type": rtype,
                "CR": cr,
                "CR Count": cr_row.get("cr_count") or len(jiras),
                "CR Title": cr_row.get("cr_title") or "",
                "CR Status": cr_row.get("cr_status") or "",
                "CR Image": cr_row.get("cr_image") or "",
                "CR Area": cr_row.get("cr_area") or "",
                "CR Subsystem": cr_row.get("cr_subsystem") or "",
                "CR Function": cr_row.get("cr_function") or "",
                "JIRA": jira.get("key") or "",
                "JIRA Title": jira.get("title") or jira.get("summary") or "",
                "JIRA Status": jira.get("status") or "",
                "JIRA Resolution": jira.get("resolution") or "",
                "Final Ticket": jira.get("final_key") or trav.get("final_key") or "",
                "Final Status": jira.get("final_status") or trav.get("final_status") or "",
                "Final Resolution": jira.get("final_resolution") or trav.get("final_resolution") or "",
                "Mapping Type": trav.get("mapping_type") or "",
                "Resolution Notes": jira.get("resolution_notes_text") or trav.get("resolution_notes_text") or "",
                "Created": jira.get("created") or "",
                "Serial No": jira.get("serial_no") or "",
                "Matched Build": jira.get("matched_build") or "",
            })
    if not rows:
        for jira in (report.get("jiras") or []):
            trav = jira.get("traversal") or {}
            info = jira.get("cr_info") or {}
            cr = trav.get("final_cr") or jira.get("cr_mapped") or "NO_CR"
            rows.append({
                "Row Type": _row_type(cr),
                "CR": cr,
                "CR Title": info.get("cr_title") or "",
                "CR Status": info.get("cr_status") or "",
                "CR Area": info.get("cr_area") or "",
                                "JIRA": jira.get("key") or "",
                "JIRA Title": jira.get("summary") or jira.get("title") or "",
                "JIRA Status": jira.get("status") or "",
                "JIRA Resolution": jira.get("resolution") or "",
                "Final Ticket": trav.get("final_key") or jira.get("final_key") or "",
                "Final Status": trav.get("final_status") or jira.get("final_status") or "",
                "Final Resolution": trav.get("final_resolution") or jira.get("final_resolution") or "",
                "Mapping Type": trav.get("mapping_type") or "",
                "Resolution Notes": jira.get("resolution_notes_text") or trav.get("resolution_notes_text") or "",
                "Matched Build": jira.get("matched_build") or "",
            })
    return rows

def _wbc_build_report_from_jql(target: Dict[str, str], build_id: str, force: bool = False) -> Dict[str, Any]:
    build = _build_tail(build_id)
    target_key = target.get("key") or target.get("name") or "target"
    cache_path = _running_build_jql_cache_path(target_key, build)
    cached = _read_json(cache_path, {})
    ttl = timedelta(minutes=30)
    now = datetime.utcnow()
    generated_at = _parse_iso_dt(cached.get("generated_at"))
    if cached.get("report") and not force and generated_at != datetime.min and now - generated_at < ttl:
        out = cached.get("report") or {}
        out["cache_status"] = "cached"
        out["cache_ttl_minutes"] = 30
        out["next_auto_refresh_at"] = (generated_at + ttl).isoformat() + "Z"
        return out
    jql = _wbc_build_jql(build)
    try:
        import sys as _sys
        scripts_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "scripts")
        if scripts_dir not in _sys.path:
            _sys.path.insert(0, scripts_dir)
        from fetch_consolidated_report import run_consolidated_report
        raw_report = run_consolidated_report(
            build_ids=[build],
            filter_id=JIRA_PDT_FILTER_ID,
            traverse=True,
            enrich_orbit=True,
            target_name=target.get("key") or target.get("name") or None,
            custom_jql=jql,
        )
        flat_rows = _wbc_flatten_consolidated_report(raw_report)
        valid_rows   = [r for r in flat_rows if not _is_invalid_row(r)]
        crs          = {str(r.get("CR") or "").strip() for r in flat_rows if _is_true_cr(str(r.get("CR") or ""))}
        mapped_jiras = {str(r.get("CR") or "").strip() for r in flat_rows if _is_cr_equiv(str(r.get("CR") or ""))}
        open_jiras   = {str(r.get("JIRA") or "").strip() for r in flat_rows if str(r.get("Row Type") or "") == "open" and str(r.get("JIRA") or "").strip()}
        all_jiras    = {str(r.get("JIRA") or "").strip() for r in flat_rows if str(r.get("JIRA") or "").strip()}
        # Valid-only counts (exclude invalid/withdrawn/won't-fix rows) — used for hero cards
        valid_crs          = {str(r.get("CR") or "").strip() for r in valid_rows if _is_true_cr(str(r.get("CR") or ""))}
        valid_mapped_jiras = {str(r.get("CR") or "").strip() for r in valid_rows if _is_cr_equiv(str(r.get("CR") or ""))}
        valid_open_jiras   = {str(r.get("JIRA") or "").strip() for r in valid_rows if str(r.get("Row Type") or "") == "open" and str(r.get("JIRA") or "").strip()}
        valid_all_jiras    = {str(r.get("JIRA") or "").strip() for r in valid_rows if str(r.get("JIRA") or "").strip()}
        report = {
            "ok": True,
            "build_id": build,
            "generated_at": now.isoformat() + "Z",
            "cache_status": "generated" if not force else "force_generated",
            "cache_ttl_minutes": 30,
            "next_auto_refresh_at": (now + ttl).isoformat() + "Z",
            "source": "JIRA JQL consolidated report",
            "jql": jql,
            "cr_count": len(crs),
            "mapped_jira_count": len(mapped_jiras),
            "open_jira_count": len(open_jiras),
            "jira_count": len(all_jiras),
            "row_count": len(flat_rows),
            "valid_cr_count": len(valid_crs),
            "valid_mapped_jira_count": len(valid_mapped_jiras),
            "valid_open_jira_count": len(valid_open_jiras),
            "valid_jira_count": len(valid_all_jiras),
            "invalid_count": len(flat_rows) - len(valid_rows),
            "rows": flat_rows,
            "summary": raw_report.get("summary") or {},
            "meta": raw_report.get("meta") or {},
        }
        _write_json(cache_path, {"generated_at": report["generated_at"], "report": report})
        return report
    except Exception as exc:
        return {
            "ok": False,
            "build_id": build,
            "generated_at": now.isoformat() + "Z",
            "cache_status": "error",
            "source": "JIRA JQL consolidated report",
            "jql": jql,
            "error": str(exc),
            "rows": [],
            "cr_count": 0,
            "jira_count": 0,
            "row_count": 0,
        }


def _running_build_report(target: Dict[str, str], db_cfg: Dict[str, str], current: Dict[str, Any], build_summary: Dict[str, Any], force: bool = False) -> Dict[str, Any]:
    """Generate WBC current-running build report from JIRA JQL, not DB tables.

    DB/config is only used upstream to discover the current Axiom running builds.
    Crash/JIRA/CR counts and detail rows come from the consolidated JQL report,
    exactly like the single-build Force Run endpoint.
    """
    running_rows = current.get("rows") or []
    running_builds = [str(r.get("build_id") or "").strip() for r in running_rows if str(r.get("build_id") or "").strip()]
    cache_key_payload = {
        "mode": "jql_consolidated_v2_qstability",
        "target": target.get("key") or target.get("name") or "",
        "running_builds": sorted(_norm_build_key(b) for b in running_builds),
        "axiom_updated_at": current.get("updated_at") or "",
        "jqls": {b: _wbc_build_jql(b) for b in running_builds},
    }
    cache_key = hashlib.sha1(json.dumps(cache_key_payload, sort_keys=True, default=str).encode("utf-8")).hexdigest()
    cache_path = _running_report_cache_path(target.get("key") or target.get("name") or "target")
    cached = _read_json(cache_path, {})
    ttl = timedelta(minutes=30)
    now = datetime.utcnow()
    saved_at = _parse_iso_dt(cached.get("saved_at") or (cached.get("report") or {}).get("generated_at"))
    if (
        not force
        and cached.get("cache_key") == cache_key
        and cached.get("report")
        and saved_at != datetime.min
        and now - saved_at < ttl
    ):
        report = cached.get("report") or {}
        report["cache_status"] = "cached"
        report["cache_key"] = cache_key
        report["cache_ttl_minutes"] = 30
        report["next_auto_refresh_at"] = (saved_at + ttl).isoformat() + "Z"
        return report

    rows_by_build: Dict[str, List[Dict[str, Any]]] = {}
    builds: List[Dict[str, Any]] = []
    errors: Dict[str, str] = {}
    missing: List[str] = []
    jql_by_build: Dict[str, str] = {}
    for running_build in running_builds:
        one = _wbc_build_report_from_jql(target, running_build, force=force)
        rows = one.get("rows") or []
        jql_by_build[running_build] = one.get("jql") or _wbc_build_jql(running_build)
        rows_by_build[running_build] = rows
        if one.get("error"):
            errors[running_build] = str(one.get("error") or "")
        if not rows:
            missing.append(running_build)
        builds.append({
            "build_id": running_build,
            "meta_id": _meta_label(running_build),
            "matched_old_builds": [],
            "cr_count": one.get("cr_count") or 0,
            "jira_count": one.get("jira_count") or 0,
            "row_count": one.get("row_count") or len(rows),
            "cache_status": one.get("cache_status") or "",
            "jql": jql_by_build[running_build],
            "source": one.get("source") or "JIRA JQL consolidated report",
        })

    report = {
        "ok": True,
        "generated_at": now.isoformat() + "Z",
        "cache_status": "generated" if not force else "force_generated",
        "cache_key": cache_key,
        "cache_ttl_minutes": 30,
        "next_auto_refresh_at": (now + ttl).isoformat() + "Z",
        "source": "JIRA JQL consolidated report",
        "running_builds": running_builds,
        "builds": builds,
        "rows_by_build": rows_by_build,
        "missing_builds": missing,
        "errors": errors,
        "jql_by_build": jql_by_build,
        "message": "Auto-generated from JIRA JQL consolidated reports; DB tables are not used for crash counts/details. Cache is reused for 30 minutes unless Force Run is clicked.",
    }
    _write_json(cache_path, {"cache_key": cache_key, "report": report, "saved_at": report["generated_at"]})
    return report


def _empty_running_build_report(current: Dict[str, Any]) -> Dict[str, Any]:
    running_builds = [str(r.get("build_id") or "").strip() for r in (current.get("rows") or []) if str(r.get("build_id") or "").strip()]
    return {
        "ok": True,
        "generated_at": "",
        "cache_status": "not_loaded",
        "source": "Current Running Builds",
        "running_builds": running_builds,
        "builds": [],
        "rows_by_build": {},
        "missing_builds": [],
        "errors": {},
        "jql_by_build": {},
        "message": "Select a current running build to generate its consolidated report.",
    }


def _target_payload(target_key: str, force_running_report: bool = False) -> Dict[str, Any]:
    targets = _discover_targets()
    target = next((t for t in targets if t["key"].lower() == str(target_key or "").lower()), None)
    if not target:
        return {"ok": False, "error": "WBC target workbook not found", "available_targets": targets}
    cfg = _load_config()
    db_cfg = (cfg.get("targets") or {}).get(target["key"], {})
    data = _load_or_sync_mainline_mtbf(target, db_cfg)
    chart_rows = data.get("chart_rows") or []
    hours = round(sum(_safe_float(r.get("hours")) for r in chart_rows), 2)
    crashes = sum(_safe_int(r.get("total_crashes")) for r in chart_rows)
    current = {"rows": [], "updated_at": "", "source": "saved_jql_tabs"}

    unique_table = db_cfg.get("unique_crs_table") or db_cfg.get("overall_crs_table") or ""
    build_summary = _build_summary_from_jiras(db_cfg.get("jiras_table") or db_cfg.get("target_table") or "", db_cfg.get("openjiras_table") or "")
    # Current Running Builds is driven by the saved JQL cards in the UI.
    # Do not use Axiom for WBC dashboard summary/current-meta values.
    running_build_report = _empty_running_build_report(current)

    return {

        "ok": True,
        "target": target,
        "db_config": db_cfg,
        "excel": data,
        "current_builds": current.get("rows") or [],
        "axiom_updated_at": current.get("updated_at") or "",
        "source": current.get("source") or "",
        "counts": {
            "rows": len(data.get("rows") or []),
            "chart_rows": len(chart_rows),
            "hours": hours,
            "crashes": crashes,
            "mtbf": round(hours / crashes, 2) if crashes else hours,
            "running_builds": len(current.get("rows") or []),
            "total_jiras": _count_from_table(db_cfg.get("jiras_table") or "", ["stability_ticket", "jira_id", "ticket"]),
            "open_jiras": _count_from_table(db_cfg.get("openjiras_table") or "", ["stability_ticket", "jira_id", "ticket"]),
            "total_crs": _count_from_table(unique_table, ["mapped_cr", "mapped_crs", "cr", "crid", "stability_ticket"]),
        },
        "previews": {
            "jiras": _preview_rows_filtered(db_cfg.get("jiras_table") or "", 100),
            "open_jiras": _preview_rows_filtered(db_cfg.get("openjiras_table") or "", 100),
            "open_crs": _preview_rows_filtered(unique_table, 100, open_cr_only=True),
            "all_crs": _preview_rows_filtered(unique_table, 150),
            "crs": _preview_rows_filtered(unique_table, 150),
                },
        "build_summary": build_summary,
        "running_build_report": running_build_report,

        "overview_summary": _load_overview_summary(target["key"]),
    }


def render_wbc_live_view_stats_page():
    return render_template("wbc_live_view_stats.html", can_edit=_can_edit())


@wbc_live_view_stats_bp.route("/wbc/live_view_status")
@login_required
def wbc_live_view_status_page():
    return render_wbc_live_view_stats_page()


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/targets")
@login_required
def api_wbc_targets():
    cfg = _load_config()
    targets = _discover_targets()
    all_targets = _discover_targets(include_hidden=True)
    for row in targets:
        row["configured"] = bool((cfg.get("targets") or {}).get(row["key"]))
        row["synced"] = os.path.exists(_target_json_path(row["key"]))
    for row in all_targets:
        row["configured"] = bool((cfg.get("targets") or {}).get(row["key"]))
        row["synced"] = os.path.exists(_target_json_path(row["key"]))
    return jsonify({"ok": True, "targets": targets, "all_targets": all_targets, "config": cfg, "can_edit": _can_edit()})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/config", methods=["GET", "POST"])
@login_required
def api_wbc_config():
    if request.method == "POST":
        if not _can_edit():
            return jsonify({"ok": False, "error": "Access denied"}), 403
        return jsonify({"ok": True, "config": _save_config(request.get_json(force=True, silent=True) or {})})
    return jsonify({"ok": True, "config": _load_config(), "targets": _discover_targets(), "can_edit": _can_edit()})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/db_tables")
@login_required
def api_wbc_db_tables():
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    return jsonify({"ok": True, "schema": _WBC_SCHEMA, "tables": _db_table_options()})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/hide", methods=["POST"])
@login_required
def api_wbc_hide_target(target_key: str):
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    cfg = _load_config()
    hidden = {_slug(x) for x in (cfg.get("hidden_targets") or [])}
    hidden.add(_slug(target_key))
    cfg["hidden_targets"] = sorted(hidden)
    cfg["updated_at"] = datetime.utcnow().isoformat() + "Z"
    _write_json(_config_path(), cfg)
    return jsonify({"ok": True, "config": cfg, "targets": _discover_targets()})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/remove", methods=["POST"])
@login_required
def api_wbc_remove_target(target_key: str):
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    key = _slug(target_key)
    cfg = _load_config()
    targets = cfg.get("targets") or {}
    targets.pop(key, None)
    cfg["targets"] = targets
    hidden = {_slug(x) for x in (cfg.get("hidden_targets") or [])}
    hidden.add(key)
    cfg["hidden_targets"] = sorted(hidden)
    cfg["updated_at"] = datetime.utcnow().isoformat() + "Z"
    _write_json(_config_path(), cfg)
    for path in (_target_json_path(key), _mtbf_json_path(key)):
        try:
            if os.path.exists(path):
                os.remove(path)
        except Exception:
            pass
    return jsonify({"ok": True, "config": cfg, "targets": _discover_targets()})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/sync", methods=["POST"])
@login_required
def api_wbc_sync():
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    payload = request.get_json(force=True, silent=True) or {}
    wanted = str(payload.get("target") or "").strip().lower()
    targets = [t for t in _discover_targets() if not wanted or t["key"].lower() == wanted]
    synced, errors = [], []
    for target in targets:
        try:
            cfg = _load_config()
            excel_path = _find_target_excel(target, (cfg.get("targets") or {}).get(target["key"], {}))
            if not excel_path or not os.path.exists(excel_path):
                raise ValueError(f"Excel file not found for {target.get('label') or target['key']}")
            sheets = _workbook_sheets(excel_path)
            if not sheets:
                raise ValueError("Workbook has no sheets")
            preferred = next((s for s in sheets if str(s).strip().lower() == "mainline_build_details"), "")
            if not preferred:
                preferred = next((s for s in sheets if "mainline" in str(s).lower() and "build" in str(s).lower()), "")
            if not preferred:
                raise ValueError("Sheet Mainline_Build_Details not found")
            data = _coerce_wbc_mtbf_payload(_sheet_to_payload(target["key"], excel_path, preferred))
            data["target_key"] = target["key"]
            data["target_label"] = target["label"]
            data["available_sheets"] = sheets
            data["mtbf_sheet"] = preferred
            data["one_time_synced"] = True
            _write_json(_mtbf_json_path(target["key"]), data)
            _write_json(_target_json_path(target["key"]), data)
            _sync_to_dashboard_mtbf_json(target["key"], data, target=target)
            synced.append({"target": target, "sheet": preferred, "rows": len(data.get("rows") or []), "chart_rows": len(data.get("chart_rows") or [])})
        except Exception as exc:
            errors.append({"target": target, "error": str(exc)})
    return jsonify({"ok": not errors or bool(synced), "synced": synced, "errors": errors, "updated_at": datetime.utcnow().isoformat() + "Z"})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>")
@login_required
def api_wbc_target(target_key: str):
    payload = _target_payload(target_key)
    return jsonify(payload), (200 if payload.get("ok") else 404)


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/saved_jql_tabs", methods=["GET"])
@login_required
def api_wbc_saved_jql_tabs(target_key: str):
    target = _find_target(target_key)
    if not target:
        return jsonify({"ok": False, "error": "WBC target not found"}), 404
    from live_view_saved_jql_service import get_cached_report_raw, list_tabs
    tabs = []

    saved_domain = _wbc_saved_jql_domain(target)
    for tab in list_tabs(_wbc_pdt_key(target), saved_domain):
        row = dict(tab)
        resolved_jql, filter_id, resolved, err = _wbc_resolve_saved_jql(row.get("jql"))
        cached = get_cached_report_raw(_wbc_pdt_key(target), saved_domain, row.get("id")) or {}
        row["raw_jql"] = row.get("jql") or ""
        row["resolved_jql"] = resolved_jql
        row["filter_id"] = filter_id
        row["filter_resolved"] = resolved
        row["filter_error"] = err
        row["build_id"] = _wbc_extract_build_id_from_jql(resolved_jql or row.get("jql") or row.get("name")) or row.get("name") or ""
        row.update(_wbc_saved_jql_cache_meta(cached))
        tabs.append(row if _can_edit() else _wbc_external_saved_jql_row(row))
    return jsonify({"ok": True, "target": target, "domain": saved_domain, "tabs": tabs})



@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/saved_jql_tabs", methods=["POST"])
@login_required
def api_wbc_saved_jql_tabs_save(target_key: str):
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    target = _find_target(target_key)
    if not target:
        return jsonify({"ok": False, "error": "WBC target not found"}), 404
    payload = request.get_json(force=True, silent=True) or {}
    try:
        from live_view_saved_jql_service import list_tabs, save_tab
        username = str(getattr(current_user, "id", "") or getattr(current_user, "username", "") or "unknown")
        pdt_key = _wbc_pdt_key(target)
        tab = save_tab(
            pdt_key,
            _wbc_saved_jql_domain(target),
            tab_id=str(payload.get("id") or "").strip() or None,
            name=str(payload.get("name") or "").strip(),
            jql=str(payload.get("jql") or "").strip(),
            username=username,
        )
        return jsonify({"ok": True, "tab": tab, "tabs": list_tabs(pdt_key, _wbc_saved_jql_domain(target))})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/saved_jql_tabs/<tab_id>", methods=["DELETE"])
@login_required
def api_wbc_saved_jql_tabs_delete(target_key: str, tab_id: str):
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    target = _find_target(target_key)
    if not target:
        return jsonify({"ok": False, "error": "WBC target not found"}), 404
    from live_view_saved_jql_service import delete_tab, list_tabs
    pdt_key = _wbc_pdt_key(target)
    saved_domain = _wbc_saved_jql_domain(target)
    deleted = delete_tab(pdt_key, saved_domain, tab_id)
    return jsonify({"ok": True, "deleted": bool(deleted), "tabs": list_tabs(pdt_key, saved_domain)})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/saved_jql_tabs/<tab_id>/report", methods=["GET", "POST"])
@login_required
def api_wbc_saved_jql_tab_report(target_key: str, tab_id: str):
    external_viewer = not _can_edit()
    # View-only users may read the same consolidated report tables, but not
    # saved-filter/JQL configuration and must never force a new run.
    if external_viewer and request.method == "POST":
        return jsonify({"ok": False, "error": "Internal target-group access required"}), 403
    target = _find_target(target_key)
    if not target:
        return jsonify({"ok": False, "error": "WBC target not found"}), 404
    force = str(request.args.get("force") or "").lower() in ("1", "true", "yes", "y")
    from live_view_saved_jql_service import get_cached_report, get_tab, set_cached_report
    saved_domain = _wbc_saved_jql_domain(target)
    tab = get_tab(_wbc_pdt_key(target), saved_domain, tab_id)
    if not tab:
        return jsonify({"ok": False, "error": "Saved JQL tab not found"}), 404
    raw_jql = str(tab.get("jql") or "").strip()
    jql, filter_id, resolved, resolve_error = _wbc_resolve_saved_jql(raw_jql)
    if not jql:
        return jsonify({"ok": False, "error": "Saved JQL is empty", "tab": tab}), 400
    if not force:
        cached = get_cached_report(_wbc_pdt_key(target), saved_domain, tab_id)
        cached_resolved_jql = str((cached or {}).get("resolved_jql") or "").strip()
        cached_effective_jql = str((cached or {}).get("jql") or "").strip()
        cached_raw_jql = str((cached or {}).get("raw_jql") or "").strip()
        cached_filter_id = str((cached or {}).get("filter_id") or "").strip()
        # For saved-filter rows, NEVER accept cache only because the raw filter
        # ID matches.  The filter ID is stable while the actual JQL/build/meta can
        # change in Jira.  Reuse cache only when the cached resolved/effective JQL
        # is exactly the same as the current filter-resolved JQL.
        if filter_id:
            jql_match = (
                cached_resolved_jql == jql
                or (
                    cached_effective_jql == jql
                    and cached_effective_jql not in (raw_jql, filter_id, f"filter = {filter_id}")
                )
            )
        else:
            jql_match = (
                cached_resolved_jql == jql
                or cached_effective_jql == jql
                or cached_raw_jql == raw_jql
            )
        if cached and jql_match:
            cached = dict(cached)
            # Flatten rows from hierarchical_report if scheduler stored raw report
            rows = cached.get("rows") or cached.get("flat_rows") or []
            if not rows and (cached.get("hierarchical_report") or cached.get("jiras")):
                rows = _wbc_flatten_consolidated_report(cached)
                cached["rows"] = rows
                cached["flat_rows"] = rows
                cached["row_count"] = len(rows)
            # Add WBC classification counts if missing (scheduler doesn't classify)
            if rows and "valid_cr_count" not in cached:
                _vrows = [r for r in rows if not _is_invalid_row(r)]
                _vcrs   = {str(r.get("CR") or "").strip() for r in _vrows if _is_true_cr(str(r.get("CR") or ""))}
                _vmapped= {str(r.get("CR") or "").strip() for r in _vrows if _is_cr_equiv(str(r.get("CR") or ""))}
                _vopen  = {str(r.get("JIRA") or "").strip() for r in _vrows if str(r.get("Row Type") or "") == "open" and str(r.get("JIRA") or "").strip()}
                _vall   = {str(r.get("JIRA") or "").strip() for r in _vrows if str(r.get("JIRA") or "").strip()}
                cached.update({
                    "valid_cr_count": len(_vcrs),
                    "valid_mapped_jira_count": len(_vmapped),
                    "valid_open_jira_count": len(_vopen),
                    "valid_jira_count": len(_vall),
                    "invalid_count": len(rows) - len(_vrows),
                    "row_count": cached.get("row_count") or len(rows),
                })
            cached.update({"ok": True, "from_cache": True, "tab": tab, "jql": jql, "raw_jql": raw_jql, "filter_id": filter_id, "filter_resolved": resolved})
            cached.update(_wbc_saved_jql_cache_meta(cached))
            return jsonify(_wbc_external_saved_jql_report(cached) if external_viewer else cached)


    now = datetime.utcnow()

    try:
        import sys as _sys
        scripts_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "scripts")
        if scripts_dir not in _sys.path:
            _sys.path.insert(0, scripts_dir)
        from fetch_consolidated_report import run_consolidated_report
        build_id = _wbc_extract_build_id_from_jql(jql) or _wbc_extract_build_id_from_jql(tab.get("name"))
        raw_report = run_consolidated_report(
            build_ids=[build_id] if build_id else [],
            filter_id=filter_id or JIRA_PDT_FILTER_ID,
            traverse=True,
            enrich_orbit=True,
            target_name=target.get("key") or target.get("name") or None,
            custom_jql=jql,
        )
        rows = _wbc_flatten_consolidated_report(raw_report)
        valid_rows   = [r for r in rows if not _is_invalid_row(r)]
        crs          = {str(r.get("CR") or "").strip() for r in rows if _is_true_cr(str(r.get("CR") or ""))}
        mapped_jiras = {str(r.get("CR") or "").strip() for r in rows if _is_cr_equiv(str(r.get("CR") or ""))}
        open_jiras   = {str(r.get("JIRA") or "").strip() for r in rows if str(r.get("Row Type") or "") == "open" and str(r.get("JIRA") or "").strip()}
        all_jiras    = {str(r.get("JIRA") or "").strip() for r in rows if str(r.get("JIRA") or "").strip()}
        # Valid-only counts (exclude invalid/withdrawn/won't-fix rows) — used for hero cards
        valid_crs          = {str(r.get("CR") or "").strip() for r in valid_rows if _is_true_cr(str(r.get("CR") or ""))}
        valid_mapped_jiras = {str(r.get("CR") or "").strip() for r in valid_rows if _is_cr_equiv(str(r.get("CR") or ""))}
        valid_open_jiras   = {str(r.get("JIRA") or "").strip() for r in valid_rows if str(r.get("Row Type") or "") == "open" and str(r.get("JIRA") or "").strip()}
        valid_all_jiras    = {str(r.get("JIRA") or "").strip() for r in valid_rows if str(r.get("JIRA") or "").strip()}
        report = {
            "ok": True,

            "tab": tab,
            "target": target,
            "generated_at": now.isoformat() + "Z",
            "cache_status": "generated" if not force else "force_generated",
            "from_cache": False,
            "source": "WBC Saved JQL consolidated report",
            "jql": jql,
            "raw_jql": raw_jql,
            "resolved_jql": jql,
            "filter_id": filter_id,
            "filter_resolved": resolved,
            "filter_error": resolve_error,
            "build_id": build_id,
            "rows": rows,
            "flat_rows": rows,
            "row_count": len(rows),
            "cr_count": len(crs),
            "mapped_jira_count": len(mapped_jiras),
            "open_jira_count": len(open_jiras),
            "jira_count": len(all_jiras),
            "valid_cr_count": len(valid_crs),
            "valid_mapped_jira_count": len(valid_mapped_jiras),
            "valid_open_jira_count": len(valid_open_jiras),
            "valid_jira_count": len(valid_all_jiras),
            "invalid_count": len(rows) - len(valid_rows),
            "summary": raw_report.get("summary") or {},
            "meta": raw_report.get("meta") or {},
        }
        stored = set_cached_report(_wbc_pdt_key(target), saved_domain, tab_id, report)
        report["generated_at"] = stored.get("generated_at") or report["generated_at"]
        report.update(_wbc_saved_jql_cache_meta(report))
        return jsonify(_wbc_external_saved_jql_report(report) if external_viewer else report)

    except Exception as exc:
        return jsonify({
            "ok": False,
            "tab": tab,
            "generated_at": now.isoformat() + "Z",
            "source": "WBC Saved JQL consolidated report",
            "jql": jql,
            "raw_jql": raw_jql,
            "resolved_jql": jql,
            "filter_id": filter_id,
            "filter_resolved": resolved,
            "filter_error": resolve_error,
            "run_error": str(exc),
            "rows": [],
            "flat_rows": [],
        }), 500





@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/running_build_report", methods=["GET", "POST"])
@login_required
def api_wbc_running_build_report(target_key: str):
    req_json = request.get_json(force=True, silent=True) if request.method == "POST" else {}
    force = str(request.args.get("force") or (req_json or {}).get("force") or "").lower() in ("1", "true", "yes", "y")
    build_id = str(request.args.get("build_id") or (req_json or {}).get("build_id") or "").strip()
    target = _find_target(target_key)
    if not target:
        return jsonify({"ok": False, "error": "WBC target not found"}), 404
    if build_id:
        one = _wbc_build_report_from_jql(target, build_id, force=force)
        rr = {
            "ok": bool(one.get("ok")),
            "generated_at": one.get("generated_at") or datetime.utcnow().isoformat() + "Z",
            "cache_status": one.get("cache_status") or "generated",
            "source": one.get("source") or "JIRA JQL consolidated report",
            "jql": one.get("jql") or "",
            "running_builds": [build_id],
            "builds": [{
                "build_id": build_id,
                "meta_id": _meta_label(build_id),
                "matched_old_builds": [],
                "cr_count": one.get("cr_count") or 0,
                "jira_count": one.get("jira_count") or 0,
                "row_count": one.get("row_count") or 0,
                "cache_status": one.get("cache_status") or "",
                "jql": one.get("jql") or "",
            }],
            "rows_by_build": {build_id: one.get("rows") or []},
            "missing_builds": [] if one.get("rows") else [build_id],
            "error": one.get("error") or "",
            "cache_ttl_minutes": one.get("cache_ttl_minutes") or 30,
            "next_auto_refresh_at": one.get("next_auto_refresh_at") or "",
            "message": "JQL-based WBC running build report. Cache is reused for 30 minutes unless Force Run is clicked.",
        }
        return jsonify({"ok": True, "target": target, "build_id": build_id, "forced": force, "running_build_report": rr})
    payload = _target_payload(target_key, force_running_report=force)
    if not payload.get("ok"):
        return jsonify(payload), 404
    return jsonify({"ok": True, "target": payload.get("target"), "build_id": build_id, "forced": force, "running_build_report": payload.get("running_build_report") or {}})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/overview_summary", methods=["GET", "POST"])
@login_required
def api_wbc_overview_summary(target_key: str):
    target = _find_target(target_key)
    if not target:
        return jsonify({"ok": False, "error": "WBC target not found"}), 404
    if request.method == "POST":
        if not _can_edit():
            return jsonify({"ok": False, "error": "Access denied"}), 403
        data = _save_overview_summary(target["key"], request.get_json(force=True, silent=True) or {})
    else:
        data = _load_overview_summary(target["key"])
    return jsonify({"ok": True, "target": target, "overview_summary": data})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/milestones", methods=["GET"])
@login_required
def api_wbc_milestones(target_key: str):
    """Fetch milestones from OneView for the given WBC target (uses PL label as SP name)."""
    target = _find_target(target_key)
    if not target:
        return jsonify({"ok": False, "error": "WBC target not found"}), 404
    # Try label first (e.g. "Kobuk.LE.3.1"), then name, then key slug
    # OneView SP names use dot notation matching the label
    label = str(target.get("label") or "").strip()
    name  = str(target.get("name")  or "").strip()
    key   = str(target.get("key")   or target_key).strip()
    # Build candidate list: label variants first, then name, then key
    candidates = []
    for raw in (label, name, key):
        if not raw:
            continue
        candidates.append(raw)
        # Also try with dots replaced by underscores and vice versa
        if "." in raw:
            candidates.append(raw.replace(".", "_"))
        if "_" in raw:
            candidates.append(raw.replace("_", "."))
    # Deduplicate preserving order
    seen_c: set = set()
    sp_candidates = [c for c in candidates if c and not (c in seen_c or seen_c.add(c))]
    # Allow caller to override the SP name via ?sp_name= query param (from the UI input)
    sp_override = str(request.args.get("sp_name") or "").strip()
    if sp_override:
        sp_candidates = [sp_override] + [c for c in sp_candidates if c != sp_override]
    sp_name = sp_candidates[0] if sp_candidates else key
    try:
        from dashboard_common import fetch_milestones_for_sp
        key_dates, source = {"ES": None, "FC": None, "CS": None, "CS1": None}, "manual"
        last_error = ""
        for candidate in sp_candidates:
            try:
                key_dates, source = fetch_milestones_for_sp(candidate)
                if any(v for v in key_dates.values()):
                    sp_name = candidate
                    break
            except Exception as _e:
                last_error = str(_e)
                continue
        if not any(v for v in key_dates.values()) and last_error:
            # Return graceful response — don't raise; let the UI show the status
            is_timeout = any(kw in last_error.lower() for kw in ("timeout", "timed out", "connection", "refused", "unreachable"))
            friendly = (
                f"OneView server unreachable (timeout). Tried SP names: {', '.join(sp_candidates[:4])}."
                if is_timeout else
                f"No milestone data found. Tried: {', '.join(sp_candidates[:4])}. Error: {last_error}"
            )
            return jsonify({"ok": False, "error": friendly, "sp_name": sp_name, "timeout": is_timeout}), 200
        # Format dates for display
        def _fmt(v):
            if not v:
                return ""
            s = str(v).strip()
            # Try to reformat YYYY-MM-DD to DD-Mon-YYYY
            try:
                from datetime import datetime as _dt
                d = _dt.strptime(s[:10], "%Y-%m-%d")
                return d.strftime("%d-%b-%Y")
            except Exception:
                return s
        es = _fmt(key_dates.get("ES"))
        fc = _fmt(key_dates.get("FC"))
        cs = _fmt(key_dates.get("CS"))
        cs1 = _fmt(key_dates.get("CS1"))
        lines = []
        if es:  lines.append(f"ES = {es}")
        if fc:  lines.append(f"FC = {fc}")
        if cs:  lines.append(f"CS = {cs}")
        if cs1 and cs1 != cs: lines.append(f"CS1 = {cs1}")
        milestone_text = "\n".join(lines) if lines else "No milestone dates found in OneView."
        return jsonify({
            "ok": True,
            "sp_name": sp_name,
            "source": source,
            "milestones": key_dates,
            "milestone_text": milestone_text,
            "es": es, "fc": fc, "cs": cs, "cs1": cs1,
        })
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc), "sp_name": sp_name}), 500


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/overview_summary/sync_excel", methods=["POST"])
@login_required
def api_wbc_overview_summary_sync_excel(target_key: str):
    """Read SWPDT_Summary (→ Summary panel) and SWPDT_Summary2 (→ PDT STATUS panel)
    from the target's MTBF workbook and save to the overview_summary JSON cache."""
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    try:
        target = _find_target(target_key)
        if not target:
            return jsonify({"ok": False, "error": "WBC target not found"}), 404
        cfg = _load_config()
        db_cfg = (cfg.get("targets") or {}).get(target["key"], {})
        excel_path = _find_target_excel(target, db_cfg)
        if not excel_path or not os.path.exists(excel_path):
            return jsonify({"ok": False, "error": "Target workbook not found"}), 404
        extracted = _read_swpdt_summary_from_excel(excel_path)
        if not extracted.get("overview") and not extracted.get("pdt_status"):
            return jsonify({"ok": False, "error": "SWPDT_Summary / SWPDT_Summary2 sheets not found in workbook"}), 404
        # Merge with existing summary (preserve engineer/title/highlights if already set)
        existing = _load_overview_summary(target["key"])
        existing["overview"] = extracted["overview"]
        existing["pdt_status"] = extracted["pdt_status"]
        existing["next_steps"] = extracted["pdt_status"]
        existing["updated_at"] = datetime.utcnow().isoformat() + "Z"
        existing["updated_by"] = str(getattr(current_user, "id", "") or getattr(current_user, "username", "") or "").strip()
        existing["source"] = "excel_swpdt"
        _write_json(_overview_summary_path(target["key"]), existing)
        return jsonify({"ok": True, "target": target, "overview_summary": existing,
                        "overview_chars": len(extracted["overview"]),
                        "pdt_status_chars": len(extracted["pdt_status"])})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 500


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/overview_summary/upload_excel", methods=["POST"])
@login_required
def api_wbc_overview_summary_upload_excel(target_key: str):
    """Upload an Excel file, show sheet picker, then read SWPDT_Summary + SWPDT_Summary2
    from the uploaded file and save to the overview_summary JSON cache.

    The selected sheet_name is used to read the Summary panel (col 2 rows 2+).
    SWPDT_Summary2 is also read automatically from the same file for the PDT STATUS panel.
    """
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    try:
        target = _find_target(target_key)
        if not target:
            return jsonify({"ok": False, "error": "WBC target not found"}), 404
        file = request.files.get("file")
        if not file or not file.filename:
            return jsonify({"ok": False, "error": "No file uploaded"}), 400
        import tempfile
        suffix = os.path.splitext(str(file.filename))[1] or ".xlsx"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name
        sheet_name = (request.form.get("sheet_name") or "").strip()
        try:
            extracted = _read_swpdt_summary_from_excel(tmp_path)
            # If a specific sheet was selected, also read it as the summary text
            # (handles cases where the user picks a non-standard sheet name)
            if sheet_name and not extracted.get("overview"):
                import openpyxl
                wb = openpyxl.load_workbook(tmp_path, data_only=True)
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    lines = []
                    for r in range(2, ws.max_row + 1):
                        # Try col 2 first (SWPDT_Summary format), fall back to col 1
                        val = str(ws.cell(r, 2).value or "").strip().replace("\xa0", " ").strip()
                        if not val:
                            val = str(ws.cell(r, 1).value or "").strip().replace("\xa0", " ").strip()
                        if val:
                            lines.append(val)
                    extracted["overview"] = "\n".join(lines)
                wb.close()
        finally:
            try:
                os.remove(tmp_path)
            except Exception:
                pass
        existing = _load_overview_summary(target["key"])
        if extracted.get("overview"):
            existing["overview"] = extracted["overview"]
        if extracted.get("pdt_status"):
            existing["pdt_status"] = extracted["pdt_status"]
            existing["next_steps"] = extracted["pdt_status"]
        existing["updated_at"] = datetime.utcnow().isoformat() + "Z"
        existing["updated_by"] = str(getattr(current_user, "id", "") or getattr(current_user, "username", "") or "").strip()
        existing["source"] = "excel_upload_swpdt"
        _write_json(_overview_summary_path(target["key"]), existing)
        return jsonify({
            "ok": True,
            "target": target,
            "overview_summary": existing,
            "overview_chars": len(extracted.get("overview") or ""),
            "pdt_status_chars": len(extracted.get("pdt_status") or ""),
            "sheet_used": sheet_name or "SWPDT_Summary",
            "uploaded_filename": str(file.filename),
        })
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 500


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/mtbf/add_build", methods=["POST"])
@login_required
def api_wbc_mtbf_add_build(target_key: str):
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    target = _find_target(target_key)
    if not target:
        return jsonify({"ok": False, "error": "WBC target not found"}), 404
    cfg = _load_config()
    db_cfg = (cfg.get("targets") or {}).get(target["key"], {})
    data = _load_or_sync_mainline_mtbf(target, db_cfg)
    rows = list(data.get("chart_rows") or [])
    row = _normalize_mtbf_row((request.get_json(force=True, silent=True) or {}).get("row") or {}, len(rows) + 1)
    if not row.get("date") or not row.get("meta_id"):
        return jsonify({"ok": False, "error": "Date and build/meta are required"}), 400
    rows.append(row)
    data = _save_mtbf_chart_rows(target, db_cfg, rows)
    return jsonify({"ok": True, "row": row, "rows": data.get("chart_rows") or [], "row_count": len(data.get("chart_rows") or []), "excel": data})


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/mtbf/save_table", methods=["POST"])
@login_required
def api_wbc_mtbf_save_table(target_key: str):
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    target = _find_target(target_key)
    if not target:
        return jsonify({"ok": False, "error": "WBC target not found"}), 404
    cfg = _load_config()
    db_cfg = (cfg.get("targets") or {}).get(target["key"], {})
    rows = (request.get_json(force=True, silent=True) or {}).get("rows") or []
    data = _save_mtbf_chart_rows(target, db_cfg, rows)
    return jsonify({"ok": True, "rows": data.get("chart_rows") or [], "row_count": len(data.get("chart_rows") or []), "excel": data})



# ---------------------------------------------------------------------------
# PPT EXPORT - same design as WBC_Report.py build_ppt()
# ---------------------------------------------------------------------------

def _wbc_build_ppt(target_key: str):
    """Generate a PowerPoint for the given WBC target from the same data used by the UI."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        raise RuntimeError("python-pptx is not installed")

    BLUE = RGBColor(0x1f, 0x5f, 0x91)
    BLUE_DARK = RGBColor(0x1a, 0x4f, 0x7b)
    WHITE = RGBColor(0xff, 0xff, 0xff)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    TEAL = RGBColor(0x15, 0x60, 0x82)
    LIGHT = RGBColor(0xee, 0xf5, 0xff)
    ROW = RGBColor(0xe9, 0xed, 0xf3)
    ROW_ALT = RGBColor(0xf7, 0xf2, 0xf6)
    BORDER = RGBColor(0xd6, 0xde, 0xea)
    GOLD = RGBColor(0xc7, 0x8b, 0x12)
    RED = RGBColor(0xd9, 0x30, 0x25)

    def _in(value):
        return Inches(float(value))

    def _plain(value, max_chars=None):
        text = re.sub(r"<[^>]+>", " ", str(value or ""))
        text = re.sub(r"&nbsp;", " ", text)
        text = re.sub(r"&amp;", "&", text)
        text = re.sub(r"&lt;", "<", text)
        text = re.sub(r"&gt;", ">", text)
        text = re.sub(r"\s+", " ", text).strip()
        if max_chars and len(text) > max_chars:
            return text[: max_chars - 1].rstrip() + "…"
        return text

    def _fmt(value):
        if isinstance(value, float):
            return f"{value:,.2f}".rstrip("0").rstrip(".")
        return str(value if value not in (None, "") else "-")

    def _add_slide(prs):
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _rect(slide, x, y, w, h, fc, lc=None, lw=0.5):
        shape = slide.shapes.add_shape(1, _in(x), _in(y), _in(w), _in(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fc
        if lc:
            shape.line.color.rgb = lc
            shape.line.width = Pt(lw)
        else:
            shape.line.fill.background()
        return shape

    def _text(slide, x, y, w, h, text, size=8, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = _plain(text)
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        return box

    def _cell(cell, text, size=6, bold=False, color=BLACK, bg=None, align=PP_ALIGN.CENTER):
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        cell.text = ""
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(1.5)
        tf.margin_right = Pt(1.5)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        try:
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = _plain(text, 260)
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    def _table(slide, x, y, w, h, headers, rows, widths=None, font_size=5.8, title_cols=None, max_text=180):
        title_cols = set(title_cols or [])
        rows = rows or []
        n_cols = max(1, len(headers))
        n_rows = max(1, len(rows)) + 1
        shape = slide.shapes.add_table(n_rows, n_cols, _in(x), _in(y), _in(w), _in(h))
        tbl = shape.table
        widths = widths or [1] * n_cols
        total = float(sum(widths) or 1)
        for i, width in enumerate(widths[:n_cols]):
            tbl.columns[i].width = int(_in(w) * (width / total))
        for ci, header in enumerate(headers):
            _cell(tbl.cell(0, ci), header, size=max(font_size, 5.2), bold=True, color=WHITE, bg=BLUE)
        if rows:
            for ri, row in enumerate(rows, start=1):
                bg = ROW if ri % 2 else ROW_ALT
                for ci in range(n_cols):
                    value = row[ci] if isinstance(row, (list, tuple)) and ci < len(row) else ""
                    _cell(tbl.cell(ri, ci), _plain(value, max_text), size=font_size, bg=bg,
                          align=PP_ALIGN.LEFT if ci in title_cols else PP_ALIGN.CENTER)
        else:
            for ci in range(n_cols):
                _cell(tbl.cell(1, ci), "No Data" if ci == 0 else "", size=font_size, bg=ROW)
        return tbl

    def _section(slide, x, y, title, w=2.4):
        return _text(slide, x, y, w, 0.2, title, size=8.2, bold=True, color=TEAL)

    def _kpi(slide, x, y, w, h, label, value):
        _rect(slide, x, y, w, h, BLUE, BLUE)
        _text(slide, x + 0.03, y + 0.05, w - 0.06, 0.16, label, size=6.6, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(slide, x + 0.03, y + 0.29, w - 0.06, 0.2, value, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    def _preview_table_rows(preview, preferred, limit):
        preview = preview or {}
        columns = list(preview.get("columns") or [])
        rows = list(preview.get("rows") or [])[:limit]
        if not columns and rows:
            columns = list(rows[0].keys())
        selected = []
        normalized = {_norm(c): c for c in columns}
        for aliases in preferred:
            found = next((normalized.get(_norm(a)) for a in aliases if normalized.get(_norm(a))), "")
            if found and found not in selected:
                selected.append(found)
        for col in columns:
            if col not in selected and len(selected) < len(preferred):
                selected.append(col)
        return selected, [[r.get(c, "") for c in selected] for r in rows]

    def _current_jql_summary(target_key):
        tabs, report_rows = [], []
        try:
            from live_view_saved_jql_service import get_cached_report, list_tabs
            for tab in list_tabs(target_key, _wbc_saved_jql_domain()):
                row = dict(tab)
                resolved_jql, filter_id, resolved, err = _wbc_resolve_saved_jql(row.get("jql"))
                build_id = _wbc_extract_build_id_from_jql(resolved_jql or row.get("jql") or row.get("name")) or row.get("name") or "-"
                cached = get_cached_report(target_key, _wbc_saved_jql_domain(), row.get("id")) or {}
                cached_rows = cached.get("rows") or cached.get("flat_rows") or []
                tabs.append({
                    "build_id": build_id,
                    "jql": resolved_jql or row.get("jql") or "",
                    "row_count": cached.get("row_count", len(cached_rows)) if cached else "-",
                    "cr_count": cached.get("cr_count", "-") if cached else "-",
                    "jira_count": cached.get("jira_count", "-") if cached else "-",
                    "generated_at": str(cached.get("generated_at") or "")[:19],
                })
                if cached_rows:
                    report_rows.extend(cached_rows[:25])
        except Exception:
            pass
        return tabs, report_rows

    target = _find_target(target_key)
    if not target:
        raise ValueError(f"WBC target not found: {target_key}")
    cfg = _load_config()
    db_cfg = (cfg.get("targets") or {}).get(target["key"], {})
    data = _target_payload(target["key"])
    if not data.get("ok"):
        raise ValueError(data.get("error") or f"WBC target not found: {target_key}")

    project = target.get("label") or target.get("name") or target_key
    chart_rows = (data.get("excel") or {}).get("chart_rows") or []
    overview = data.get("overview_summary") or {}
    counts = data.get("counts") or {}
    previews = data.get("previews") or {}
    jql_tabs, cached_report_rows = _current_jql_summary(target["key"])
    last_row = chart_rows[-1] if chart_rows else {}
    first_jql = jql_tabs[0] if jql_tabs else {}
    current_meta = first_jql.get("build_id") or last_row.get("crm_build_id") or last_row.get("meta_id") or "-"
    current_mtbf = _fmt(last_row.get("mtbf") or counts.get("mtbf"))
    current_crashes = _fmt(first_jql.get("cr_count") if first_jql else (last_row.get("crash") or last_row.get("total_crashes")))
    current_hours = _fmt(last_row.get("hours"))
    report_date = str(last_row.get("date") or datetime.now().date())[:10]
    open_cr_count = (previews.get("open_crs") or {}).get("count") or 0

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: portal-like overview/status slide.
    slide = _add_slide(prs)
    _text(slide, 0.16, 0.08, 5.9, 0.28, f"Current Meta: {current_meta}", size=12, bold=True)
    _text(slide, 6.48, 0.05, 6.55, 0.15, f"PDT WBC Stability Dashboard : {project}", size=5.0)
    _table(slide, 0.20, 0.58, 6.00, 0.78,
           [f"Date: {report_date}", f"{project} PDT Status", ""],
           [["Target", "OEM", "Project Timelines"], [project, "-", _plain(overview.get("pdt_status") or overview.get("next_steps") or "-", 95)]],
           widths=[1.6, 1.0, 3.4], font_size=5.8, max_text=110)
    kpis = [
        ("Current PDT\nMTBF", current_mtbf),
        ("Current Running\nMeta", current_meta),
        ("Current META\nCrashes", current_crashes),
        ("Current META\nHours", current_hours),
        ("Running\nBuilds", str(len(jql_tabs))),
        ("Open Jira\nCurrent Meta", _fmt(counts.get("open_jiras"))),
        ("Open CR\nCurrent Meta", _fmt(open_cr_count)),
        ("Total\nCRs", _fmt(counts.get("total_crs"))),
    ]
    for idx, (label, value) in enumerate(kpis[:8]):
        _kpi(slide, 0.20 + (idx % 4) * 1.51, 1.55 + (idx // 4) * 0.62, 1.44, 0.52, label, _plain(value, 28))

    _section(slide, 0.30, 2.85, "Key Updates", w=1.2)
    _text(slide, 0.32, 3.15, 5.82, 1.02, overview.get("overview") or overview.get("summary_title") or "No summary entered.", size=6.7)
    _section(slide, 0.30, 4.48, "MTBF Chart", w=1.2)
    recent_mtbf = chart_rows[-3:] if chart_rows else []
    _table(slide, 0.32, 4.86, 5.88, 0.68,
           ["Team", "Meta", "Total Hours", "Total Crashes", "MTBF"],
           [["PDT", r.get("crm_build_id") or r.get("meta_id"), _fmt(r.get("hours")), _fmt(r.get("crash") or r.get("total_crashes")), _fmt(r.get("mtbf"))] for r in recent_mtbf],
           widths=[0.7, 1.8, 1.05, 1.05, 0.75], font_size=5.8, max_text=80)
    _section(slide, 0.30, 5.78, "Weekly Stability Stats (SW PDT)", w=2.6)
    _table(slide, 0.32, 6.14, 5.90, 0.82,
           ["Team", "Builds Tested", "Total Hours", "Total Crashes"],
           [["PDT", str(len(chart_rows)), _fmt(counts.get("hours")), _fmt(counts.get("crashes"))]],
           widths=[1.0, 2.1, 1.5, 1.4], font_size=5.7)
    _rect(slide, 6.38, 0.15, 0.01, 7.12, BORDER, BORDER)
    _section(slide, 6.50, 0.50, "Saved JQL / Current Running Builds", w=2.8)
    _table(slide, 6.55, 0.86, 6.45, 1.55,
           ["S.No.", "Build ID", "CRs", "JIRAs", "Rows", "Last Run"],
           [[str(i), r.get("build_id"), _fmt(r.get("cr_count")), _fmt(r.get("jira_count")), _fmt(r.get("row_count")), r.get("generated_at") or "-"] for i, r in enumerate(jql_tabs[:5], start=1)],
           widths=[0.45, 2.6, 0.55, 0.55, 0.55, 1.75], font_size=4.9, title_cols={1}, max_text=80)
    _section(slide, 6.50, 2.70, "Open CR Details", w=1.6)
    cr_cols, cr_rows = _preview_table_rows(previews.get("open_crs"), [["mapped_cr", "cr", "cr_id"], ["cr_title", "title", "summary"], ["cr_area", "area"], ["cr_status", "status"]], 5)
    _table(slide, 6.55, 3.04, 6.45, 2.30,
           [c.replace("_", " ") for c in cr_cols], cr_rows,
           widths=[0.95, 3.1, 1.0, 0.95][:len(cr_cols)], font_size=4.7, title_cols={1}, max_text=120)
    _section(slide, 6.50, 5.73, "Open JIRA Details", w=1.55)
    jira_cols, jira_rows = _preview_table_rows(previews.get("open_jiras"), [["stability_ticket", "jira", "jira_id"], ["jira_title", "title", "summary"], ["status", "jira_status"], ["created", "jira_date", "updated"]], 2)
    _table(slide, 6.55, 6.08, 6.45, 0.88,
           [c.replace("_", " ") for c in jira_cols], jira_rows,
           widths=[1.0, 3.3, 0.8, 0.9][:len(jira_cols)], font_size=4.5, title_cols={1}, max_text=120)

    # Slide 2: custom MTBF trend chart.
    slide = _add_slide(prs)
    _rect(slide, 0, 0, 13.33, 7.5, LIGHT)
    _text(slide, 0.22, 0.14, 3.0, 0.22, "MTBF Trend by Build", size=8, bold=True)
    _rect(slide, 0.18, 0.44, 12.95, 6.78, WHITE, RGBColor(0xe8, 0xee, 0xf6), 0.5)
    if chart_rows:
        max_points = 42
        rows = chart_rows[-max_points:]
        cats = [str(r.get("crm_build_id") or r.get("meta_id") or "") for r in rows]
        hours = [_safe_float(r.get("hours")) for r in rows]
        crashes = [_safe_float(r.get("crash") or r.get("total_crashes")) for r in rows]
        mtbf = [_safe_float(r.get("mtbf")) for r in rows]
        left, top, width, height = 0.90, 0.92, 11.35, 4.85
        bottom = top + height
        n = max(1, len(rows))
        h_max = max([1.0] + hours + crashes)
        m_max = max([1.0] + mtbf)
        h_max = (int(h_max / 500) + 1) * 500 if h_max > 500 else max(10, (int(h_max / 10) + 1) * 10)
        m_max = (int(m_max / 200) + 1) * 200 if m_max > 200 else max(10, (int(m_max / 10) + 1) * 10)
        _text(slide, left, 0.58, width, 0.20, "MTBF by Build", size=7.4, bold=True, align=PP_ALIGN.CENTER)
        for i in range(6):
            y = bottom - (height * i / 5.0)
            _rect(slide, left, y, width, 0.004, RGBColor(0xee, 0xf1, 0xf5), RGBColor(0xee, 0xf1, 0xf5))
            _text(slide, left - 0.48, y - 0.06, 0.38, 0.12, int(h_max * i / 5.0), size=4.5, color=RGBColor(0x61, 0x6f, 0x82), align=PP_ALIGN.RIGHT)
            _text(slide, left + width + 0.08, y - 0.06, 0.38, 0.12, int(m_max * i / 5.0), size=4.5, color=GOLD)
        step = width / n
        bar_w = min(0.09, step * 0.38)
        points = []
        for i, cat in enumerate(cats):
            cx = left + step * (i + 0.5)
            bh = 0 if h_max <= 0 else height * (hours[i] / h_max)
            _rect(slide, cx - bar_w / 2, bottom - bh, bar_w, max(0.015, bh), RGBColor(0x3b, 0x5b, 0xdb))
            ch = 0 if h_max <= 0 else height * (crashes[i] / h_max)
            dot = slide.shapes.add_shape(9, _in(cx - 0.025), _in(bottom - ch - 0.025), _in(0.05), _in(0.05))
            dot.fill.solid(); dot.fill.fore_color.rgb = RED; dot.line.fill.background()
            mh = 0 if m_max <= 0 else height * (mtbf[i] / m_max)
            points.append((cx, bottom - mh))
        for p1, p2 in zip(points, points[1:]):
            line = slide.shapes.add_connector(1, _in(p1[0]), _in(p1[1]), _in(p2[0]), _in(p2[1]))
            line.line.color.rgb = GOLD; line.line.width = Pt(1.05)
        for x, y in points:
            marker = slide.shapes.add_shape(9, _in(x - 0.025), _in(y - 0.025), _in(0.05), _in(0.05))
            marker.fill.solid(); marker.fill.fore_color.rgb = GOLD; marker.line.color.rgb = WHITE; marker.line.width = Pt(0.35)
        label_step = max(1, int((n + 17) / 18))
        for i, cat in enumerate(cats):
            if i % label_step != 0 and i != n - 1:
                continue
            lab = _text(slide, left + step * (i + 0.5) - 0.28, bottom + 0.08, 0.56, 0.45, _plain(cat, 28), size=3.6, color=RGBColor(0x45, 0x52, 0x63), align=PP_ALIGN.RIGHT)
            lab.rotation = 315
        _text(slide, 5.55, 6.66, 2.3, 0.2, "● Hours    ● Crashes    ● MTBF", size=5.0, color=RGBColor(0x45, 0x52, 0x63), align=PP_ALIGN.CENTER)
    else:
        _text(slide, 0.60, 1.30, 11.8, 0.5, "No MTBF chart data available.", size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Detail slides from the current UI datasets.
    detail_specs = [
        ("Open CRs", previews.get("open_crs"), [["mapped_cr", "cr", "cr_id"], ["cr_title", "title", "summary"], ["cr_area", "area"], ["cr_status", "status"], ["cr_age", "age"]], 18),
        ("Open JIRAs", previews.get("open_jiras"), [["stability_ticket", "jira", "jira_id"], ["jira_title", "title", "summary"], ["status", "jira_status"], ["created", "jira_date", "updated"]], 18),
        ("All CRs", previews.get("all_crs"), [["mapped_cr", "cr", "cr_id"], ["cr_title", "title", "summary"], ["cr_area", "area"], ["cr_status", "status"]], 18),
        ("JIRAs", previews.get("jiras"), [["stability_ticket", "jira", "jira_id"], ["jira_title", "title", "summary"], ["status", "jira_status"], ["created", "jira_date", "updated"]], 18),
    ]
    for title, preview, preferred, limit in detail_specs:
        cols, rows = _preview_table_rows(preview, preferred, limit)
        if not cols and not rows:
            continue
        slide = _add_slide(prs)
        _text(slide, 0.30, 0.28, 4.8, 0.25, f"{title} - {project}", size=11, bold=True, color=TEAL)
        _table(slide, 0.28, 0.78, 12.78, 6.15, [c.replace("_", " ") for c in cols], rows,
               font_size=4.8, title_cols={1}, max_text=150)

    if cached_report_rows:
        cols = []
        preferred = ["CR", "CR Title", "CR Status", "CR Area", "JIRA", "JIRA Title", "JIRA Status", "Final Ticket"]
        for col in preferred:
            if any(col in r for r in cached_report_rows):
                cols.append(col)
        for col in list(cached_report_rows[0].keys()):
            if col not in cols and len(cols) < 8:
                cols.append(col)
        slide = _add_slide(prs)
        _text(slide, 0.30, 0.28, 5.8, 0.25, f"Saved JQL Cached Report - {project}", size=11, bold=True, color=TEAL)
        _table(slide, 0.28, 0.78, 12.78, 6.15, cols, [[r.get(c, "") for c in cols] for r in cached_report_rows[:18]],
               font_size=4.5, title_cols={1, 5}, max_text=150)

    if chart_rows:
        slide = _add_slide(prs)
        _text(slide, 0.30, 0.28, 5.8, 0.25, f"Mainline Build Details - {project}", size=11, bold=True, color=TEAL)
        _table(slide, 0.28, 0.78, 12.78, 6.15,
               ["S.No", "CRM Build ID", "Date", "Hours+", "Crash", "MTBF"],
               [[r.get("s_no"), r.get("crm_build_id") or r.get("meta_id"), r.get("date"), _fmt(r.get("hours")), _fmt(r.get("crash") or r.get("total_crashes")), _fmt(r.get("mtbf"))] for r in chart_rows[-24:]],
               widths=[0.45, 3.4, 1.0, 1.0, 1.0, 1.0], font_size=5.4, title_cols={1}, max_text=120)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _wbc_fr_sheet_payload(excel_path: str) -> Dict[str, Any]:
    """Read the PL-wise FR_Analysis worksheet as an editable table.

    Columns whose header cell is blank are silently dropped so the UI
    does not show a sea of empty 'Column N' placeholders.
    """
    import openpyxl
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    try:
        if "FR_Analysis" not in wb.sheetnames:
            raise ValueError("FR_Analysis sheet was not found in the PL workbook.")
        ws = wb["FR_Analysis"]
        raw_headers = [str(ws.cell(1, col).value or "").strip() for col in range(1, ws.max_column + 1)]
        # Keep only columns that have a non-empty header
        non_empty_idx = [i for i, h in enumerate(raw_headers) if h]
        if not non_empty_idx:
            # Fallback: keep all and use Column N names
            non_empty_idx = list(range(len(raw_headers)))
            raw_headers = [h or f"Column {i+1}" for i, h in enumerate(raw_headers)]
        headers = [raw_headers[i] for i in non_empty_idx]
        col_nums = [i + 1 for i in non_empty_idx]  # 1-based column numbers
        rows = []
        for row_vals in ws.iter_rows(min_row=2, values_only=True):
            selected = [
                ("" if row_vals[i] is None else str(row_vals[i]).strip().replace("\xa0", " ").strip())
                if i < len(row_vals) else ""
                for i in non_empty_idx
            ]
            if any(v for v in selected):
                rows.append(selected)
        return {"sheet": ws.title, "excel_path": excel_path, "headers": headers, "rows": rows}
    finally:
        wb.close()


def _save_wbc_fr_sheet(excel_path: str, headers: List[Any], rows: List[List[Any]]) -> Dict[str, Any]:
    """Replace only FR_Analysis table values while retaining the rest of the PL workbook."""
    import openpyxl
    clean_headers = [str(value or "").strip() for value in headers if str(value or "").strip()]
    if not clean_headers:
        raise ValueError("At least one FR Analysis column is required.")
    wb = openpyxl.load_workbook(excel_path)
    try:
        ws = wb["FR_Analysis"] if "FR_Analysis" in wb.sheetnames else wb.create_sheet("FR_Analysis")
        for row in ws.iter_rows():
            for cell in row:
                cell.value = None
        for col, header in enumerate(clean_headers, start=1):
            ws.cell(1, col, header)
        saved_rows = 0
        for row_num, row in enumerate(rows or [], start=2):
            values = list(row) if isinstance(row, list) else []
            if not any(str(value or "").strip() for value in values):
                continue
            for col, _header in enumerate(clean_headers, start=1):
                ws.cell(row_num, col, values[col - 1] if col <= len(values) else "")
            saved_rows += 1
        wb.save(excel_path)
        return {"sheet": ws.title, "excel_path": excel_path, "headers": clean_headers, "saved_rows": saved_rows}
    finally:
        wb.close()


def _wbc_cr_analysis_sheet(excel_path: str, target_sheet: str, force: bool = False) -> Dict[str, Any]:
    """Run session-authorized QGenie summaries for a WBC workbook CR sheet."""
    try:
        import openpyxl
        from src.qgenie_service import get_current_qgenie_client, get_session_qgenie_highlights_model
    except ImportError as exc:
        raise RuntimeError(f"CR analysis dependency is unavailable: {exc}")

    client = get_current_qgenie_client()
    if not client:
        raise PermissionError("QGenie API key is not configured for this session.")

    sheet_names = (
        ("Open_CR_Details", "Open CR Details", "Open CRs", "Open CR")
        if target_sheet == "open_cr"
        else ("Current_Meta_CR", "Current Meta CR", "CR Details")
    )
    wb = openpyxl.load_workbook(excel_path)
    try:
        ws = next((wb[name] for name in sheet_names if name in wb.sheetnames), None)
        if ws is None:
            raise ValueError(f"Required CR sheet not found: {', '.join(sheet_names)}")

        header_row = next(
            (row for row in range(1, min(ws.max_row, 20) + 1)
             if any(str(ws.cell(row, col).value or "").strip() for col in range(1, ws.max_column + 1))),
            1,
        )
        headers = [str(ws.cell(header_row, col).value or "").strip() for col in range(1, ws.max_column + 1)]
        normalized = [_norm(h) for h in headers]
        cr_col = next((i + 1 for i, name in enumerate(normalized)
                       if name in ("cr", "cr_id", "crid") or name.startswith("cr_")), None)
        if not cr_col:
            raise ValueError("No CR/CR-ID column found in the selected worksheet.")

        analysis_col = next((i + 1 for i, name in enumerate(normalized)
                             if name in ("qgenie_analysis", "qgenieanalysis", "cr_analysis", "cranalysis")), None)
        if not analysis_col:
            analysis_col = ws.max_column + 1
            ws.cell(header_row, analysis_col, "Qgenie Analysis")
            headers.append("Qgenie Analysis")

        model = get_session_qgenie_highlights_model()
        today = datetime.now().strftime("%Y-%m-%d")
        processed = skipped = failed = 0
        for row_num in range(header_row + 1, ws.max_row + 1):
            cr_value = str(ws.cell(row_num, cr_col).value or "").strip()
            if not cr_value:
                continue
            existing = str(ws.cell(row_num, analysis_col).value or "").strip()
            if not force and existing and today in existing:
                skipped += 1
                continue

            context_parts = []
            for col_num, header in enumerate(headers[:analysis_col - 1], start=1):
                value = str(ws.cell(row_num, col_num).value or "").strip()
                if value and _norm(header) not in ("qgenie_analysis", "qgenieanalysis", "cr_analysis", "cranalysis"):
                    context_parts.append(f"{header}: {value}")
            prompt = (
                f"Provide a concise PDT engineering analysis for {cr_value}. "
                "State likely impact, current status/risk, and recommended next action. "
                "Use only the following workbook context when it is relevant:\n"
                + "\n".join(context_parts)[:2500]
            )
            try:
                response = client.chat(model=model, messages=[{"role": "user", "content": prompt}])
                text = str(getattr(response, "content", response) or "").strip()
                text = re.sub(r"\s+", " ", text).strip()
                if text:
                    ws.cell(row_num, analysis_col, f"[{today}] {text[:1800]}")
                    processed += 1
                else:
                    failed += 1
            except Exception:
                failed += 1

        if processed:
            wb.save(excel_path)
        return {
            "sheet": ws.title, "excel_path": excel_path, "processed": processed,
            "skipped": skipped, "failed": failed, "analysis_column": analysis_col,
        }
    finally:
        wb.close()


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/fr_analysis", methods=["GET", "POST"])
@login_required
def api_wbc_fr_analysis(target_key: str):
    """Read or save the PL-wise FR_Analysis data.

    GET  – returns JSON cache if present, otherwise reads from the PL workbook
           and seeds the JSON cache.
    POST – saves the submitted headers/rows to JSON only (no Excel write).
           The JSON is stored at:
             <store_dir>/fr_analysis/fr_analysis_<slug>.json
           mirroring the MTBF JSON pattern.
    """
    if request.method == "POST" and not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    try:
        target = _find_target(target_key)
        if not target:
            return jsonify({"ok": False, "error": "WBC target not found"}), 404
        key = target["key"]
        json_path = _fr_json_path(key)

        if request.method == "POST":
            body = request.get_json(force=True, silent=True) or {}
            raw_headers = body.get("headers") or []
            raw_rows = body.get("rows") or []
            clean_headers = [str(h or "").strip() for h in raw_headers if str(h or "").strip()]
            if not clean_headers:
                return jsonify({"ok": False, "error": "At least one FR Analysis column is required."}), 400
            saved_rows = 0
            clean_rows = []
            for row in raw_rows:
                values = list(row) if isinstance(row, list) else []
                if not any(str(v or "").strip() for v in values):
                    continue
                padded = [(str(values[i]) if i < len(values) else "") for i in range(len(clean_headers))]
                clean_rows.append(padded)
                saved_rows += 1
            payload = {
                "target_key": key,
                "target_label": target.get("label") or target.get("name") or key,
                "sheet": "FR_Analysis",
                "excel_path": "",
                "headers": clean_headers,
                "rows": clean_rows,
                "saved_rows": saved_rows,
                "updated_at": datetime.utcnow().isoformat() + "Z",
                "updated_by": str(getattr(current_user, "id", "") or getattr(current_user, "username", "") or "").strip(),
                "source": "json",
            }
            _write_json(json_path, payload)
            return jsonify({"ok": True, "target": target, **payload})

        # GET: try JSON cache first
        cached = _read_json(json_path, {})
        if cached.get("headers"):
            cached.setdefault("source", "json")
            cached.setdefault("json_path", json_path)
            return jsonify({"ok": True, "target": target, **cached})

        # Fall back to Excel read and seed the JSON cache
        cfg = _load_config()
        excel_path = _find_target_fr_workbook(target, (cfg.get("targets") or {}).get(key, {}))
        if not excel_path:
            # No Excel and no JSON cache — return an empty sheet so the user
            # can build the FR Analysis from scratch in the portal.
            return jsonify({
                "ok": True,
                "target": target,
                "target_key": key,
                "target_label": target.get("label") or target.get("name") or key,
                "sheet": "FR_Analysis",
                "excel_path": "",
                "headers": [],
                "rows": [],
                "saved_rows": 0,
                "source": "empty",
                "json_path": json_path,
                "updated_at": "",
                "updated_by": "",
            })
        result = _wbc_fr_sheet_payload(excel_path)
        seed = {
            "target_key": key,
            "target_label": target.get("label") or target.get("name") or key,
            "source": "excel_seeded",
            "json_path": json_path,
            **result,
        }
        try:
            _write_json(json_path, seed)
        except Exception:
            pass
        return jsonify({"ok": True, "target": target, **seed})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 500


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/fr_analysis/sync_excel", methods=["POST"])
@login_required
def api_wbc_fr_sync_excel(target_key: str):
    """Read FR_Analysis from the PL workbook and merge rows missing from the JSON cache.

    - If no JSON cache exists, seeds it from Excel (same as the GET fallback).
    - If a JSON cache exists, appends any Excel rows whose first-column key is
      absent from the cache, preserving all existing edits.
    - Returns the merged payload so the UI can refresh immediately.
    """
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    try:
        target = _find_target(target_key)
        if not target:
            return jsonify({"ok": False, "error": "WBC target not found"}), 404
        key = target["key"]
        json_path = _fr_json_path(key)
        cfg = _load_config()
        excel_path = _find_target_fr_workbook(target, (cfg.get("targets") or {}).get(key, {}))
        if not excel_path or not os.path.exists(excel_path):
            return jsonify({"ok": False, "error": "PL-wise FR workbook not found for this target"}), 404

        excel_data = _wbc_fr_sheet_payload(excel_path)
        excel_headers = excel_data.get("headers") or []
        excel_rows = excel_data.get("rows") or []

        # Always replace the cache with fresh Excel data so stale/empty
        # values from a previous manual upload are overwritten.
        fresh_payload = {
            "target_key": key,
            "target_label": target.get("label") or target.get("name") or key,
            "sheet": "FR_Analysis",
            "excel_path": excel_path,
            "headers": excel_headers,
            "rows": excel_rows,
            "saved_rows": len(excel_rows),
            "updated_at": datetime.utcnow().isoformat() + "Z",
            "updated_by": str(getattr(current_user, "id", "") or
                              getattr(current_user, "username", "") or "").strip(),
            "source": "excel_seeded",
        }
        _write_json(json_path, fresh_payload)
        return jsonify({"ok": True, "target": target, "added": len(excel_rows),
                        "total": len(excel_rows), **fresh_payload})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 500


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/excel/list_sheets", methods=["POST"])
@login_required
def api_wbc_excel_list_sheets(target_key: str):
    """List sheets from an uploaded Excel file (used for sheet selection UI)."""
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    try:
        file = request.files.get("file")
        if not file or not file.filename:
            return jsonify({"ok": False, "error": "No file uploaded"}), 400
        import tempfile
        suffix = os.path.splitext(str(file.filename))[1] or ".xlsx"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name
        try:
            sheets = _workbook_sheets(tmp_path)
        finally:
            try:
                os.remove(tmp_path)
            except Exception:
                pass
        return jsonify({"ok": True, "sheets": sheets, "filename": str(file.filename)})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 500


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/fr_analysis/upload_excel", methods=["POST"])
@login_required
def api_wbc_fr_upload_excel(target_key: str):
    """Upload an Excel file and read the FR_Analysis sheet from it (one-time import).
    Saves the result to the FR JSON cache so subsequent GETs serve from JSON.
    Accepts optional form field 'sheet_name' to override auto-detection.
    """
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    try:
        target = _find_target(target_key)
        if not target:
            return jsonify({"ok": False, "error": "WBC target not found"}), 404
        key = target["key"]
        json_path = _fr_json_path(key)
        file = request.files.get("file")
        if not file or not file.filename:
            return jsonify({"ok": False, "error": "No file uploaded"}), 400
        import tempfile
        suffix = os.path.splitext(str(file.filename))[1] or ".xlsx"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name
        sheet_name = (request.form.get("sheet_name") or "").strip()
        try:
            if sheet_name:
                raw = _sheet_to_payload(key, tmp_path, sheet_name)
                headers = raw.get("headers") or []
                rows = raw.get("rows") or []
                # Strip columns whose header is blank
                non_empty_cols = [i for i, h in enumerate(headers) if str(h or "").strip()]
                if non_empty_cols and len(non_empty_cols) < len(headers):
                    headers = [headers[i] for i in non_empty_cols]
                    rows = [[row[i] if i < len(row) else "" for i in non_empty_cols] for row in rows]
                # Strip rows that are entirely blank
                rows = [row for row in rows if any(str(v or "").strip() for v in row)]
                result = {"headers": headers, "rows": rows, "sheet": sheet_name, "saved_rows": len(rows)}
            else:
                result = _wbc_fr_sheet_payload(tmp_path)
                headers = result.get("headers") or []
                rows = result.get("rows") or []
                non_empty_cols = [i for i, h in enumerate(headers) if str(h or "").strip()]
                if non_empty_cols and len(non_empty_cols) < len(headers):
                    headers = [headers[i] for i in non_empty_cols]
                    rows = [[row[i] if i < len(row) else "" for i in non_empty_cols] for row in rows]
                rows = [row for row in rows if any(str(v or "").strip() for v in row)]
                result["headers"] = headers
                result["rows"] = rows
                result["saved_rows"] = len(rows)
        finally:
            try:
                os.remove(tmp_path)
            except Exception:
                pass
        seed = {
            "target_key": key,
            "target_label": target.get("label") or target.get("name") or key,
            "source": "excel_upload",
            "json_path": json_path,
            "uploaded_filename": str(file.filename),
            "updated_at": datetime.utcnow().isoformat() + "Z",
            "updated_by": str(getattr(current_user, "id", "") or
                              getattr(current_user, "username", "") or "").strip(),
            **result,
        }
        _write_json(json_path, seed)
        return jsonify({"ok": True, "target": target, **seed})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 500


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/mtbf/upload_excel", methods=["POST"])
@login_required
def api_wbc_mtbf_upload_excel(target_key: str):
    """Upload an Excel file and read the Mainline_Build_Details sheet (one-time import).
    Saves the result to the MTBF JSON cache so subsequent loads serve from JSON.
    """
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    try:
        target = _find_target(target_key)
        if not target:
            return jsonify({"ok": False, "error": "WBC target not found"}), 404
        key = target["key"]
        file = request.files.get("file")
        if not file or not file.filename:
            return jsonify({"ok": False, "error": "No file uploaded"}), 400
        import tempfile
        suffix = os.path.splitext(str(file.filename))[1] or ".xlsx"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name
        sheet_name = (request.form.get("sheet_name") or "").strip()
        try:
            sheets = _workbook_sheets(tmp_path)
            if sheet_name:
                sheet = sheet_name
            else:
                sheet = next((s for s in sheets if str(s).strip().lower() == "mainline_build_details"), "")
                if not sheet:
                    sheet = next((s for s in sheets if "mainline" in str(s).lower() and "build" in str(s).lower()), "")
                if not sheet:
                    raise ValueError(
                        f"Sheet Mainline_Build_Details not found. Available: {', '.join(sheets)}"
                    )
            data = _coerce_wbc_mtbf_payload(_sheet_to_payload(key, tmp_path, sheet))
        finally:
            try:
                os.remove(tmp_path)
            except Exception:
                pass
        data["target_key"] = key
        data["target_label"] = target.get("label") or target.get("name") or key
        data["mtbf_sheet"] = sheet
        data["one_time_synced"] = True
        data["uploaded_filename"] = str(file.filename)
        data["saved_json"] = _mtbf_json_path(key)
        _write_json(_mtbf_json_path(key), data)
        _write_json(_target_json_path(key), data)
        _sync_to_dashboard_mtbf_json(key, data, target=target)
        return jsonify({
            "ok": True, "target": target, "excel": data,
            "rows": len(data.get("chart_rows") or []),
        })
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 500


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/cr_analysis", methods=["POST"])
@login_required
def api_wbc_cr_analysis(target_key: str):
    """Generate QGenie CR analysis in a target's Open CR or Current Meta workbook sheet."""
    if not _can_edit():
        return jsonify({"ok": False, "error": "Access denied"}), 403
    try:
        body = request.get_json(silent=True) or {}
        target_sheet = str(body.get("target_sheet") or "open_cr").strip().lower()
        if target_sheet not in ("open_cr", "current_cr"):
            return jsonify({"ok": False, "error": "target_sheet must be open_cr or current_cr"}), 400
        target = _find_target(target_key)
        if not target:
            return jsonify({"ok": False, "error": "WBC target not found"}), 404
        cfg = _load_config()
        db_cfg = (cfg.get("targets") or {}).get(target["key"], {})
        excel_path = _find_target_fr_workbook(target, db_cfg)
        if not excel_path or not os.path.exists(excel_path):
            return jsonify({"ok": False, "error": "PL-wise WBC FR workbook was not found"}), 404
        result = _wbc_cr_analysis_sheet(excel_path, target_sheet, bool(body.get("force")))
        return jsonify({"ok": True, **result})
    except PermissionError as exc:
        return jsonify({"ok": False, "requires_config": True, "error": str(exc)}), 401
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 500


@wbc_live_view_stats_bp.route("/api/wbc_live_view_stats/target/<path:target_key>/export_ppt")
@login_required
def api_wbc_export_ppt(target_key: str):
    """Download a PowerPoint for the given WBC target."""
    try:
        buf = _wbc_build_ppt(target_key)
        target = _find_target(target_key)
        label = (target.get("label") or target.get("name") or target_key).replace(" ", "_")
        filename = f"WBC_{label}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
        return send_file(
            buf,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as exc:
        import traceback
        return jsonify({"ok": False, "error": str(exc), "trace": traceback.format_exc()}), 500
