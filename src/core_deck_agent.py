import os
import re
import shutil
import zipfile
from datetime import datetime
from typing import Any, Dict, List
from xml.etree import ElementTree as ET


def _safe_str(value: Any) -> str:
    return str(value or "").strip()


def _safe_name(value: str) -> str:
    text = _safe_str(value) or "core_deck"
    return re.sub(r"[^A-Za-z0-9._-]+", "_", text).strip("._") or "core_deck"


class CoreDeckAgent:
    """Reference-driven Core Deck PPTX helper.

    The agent intentionally uses deterministic PPTX parsing/mapping first:
    - analyze_reference() extracts slide titles, shapes, charts, tables and colors.
    - generate_pptx() preserves the reference deck and updates only detected
      meta/build/crash/MTBF/KPI text values. Everything else stays as-is.
    """

    DATA_KEYWORDS = {
        "mtbf": ("mtbf", "mean time", "failure"),
        "cr_area": ("area", "component"),
        "cr_subsystem": ("subsystem", "sub-system"),
        "top_hitters": ("top", "hitter", "cr"),
        "open_crs": ("open", "analysis", "nosir"),
        "summary": ("summary", "overview", "kpi", "dashboard"),
        "builds": ("build", "meta", "flavor"),
    }

    def __init__(self, generated_root: str) -> None:
        self.generated_root = generated_root

    def analyze_reference(self, pptx_path: str, template_name: str = "") -> Dict[str, Any]:
        try:
            from pptx import Presentation
        except Exception as exc:  # pragma: no cover - dependency/runtime guard
            raise RuntimeError(f"python-pptx is required for reference analysis: {exc}")

        prs = Presentation(pptx_path)
        slides: List[Dict[str, Any]] = []
        colors = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            shape_rows = []
            chart_count = 0
            table_count = 0
            for shape in slide.shapes:
                stype = str(getattr(shape, "shape_type", ""))
                text = ""
                try:
                    if getattr(shape, "has_text_frame", False):
                        text = "\n".join(p.text for p in shape.text_frame.paragraphs if _safe_str(p.text))
                        if text:
                            texts.append(text)
                except Exception:
                    pass
                try:
                    if getattr(shape, "has_chart", False):
                        chart_count += 1
                except Exception:
                    pass
                try:
                    if getattr(shape, "has_table", False):
                        table_count += 1
                except Exception:
                    pass
                try:
                    fill = getattr(shape, "fill", None)
                    fore = getattr(fill, "fore_color", None)
                    rgb = getattr(fore, "rgb", None)
                    if rgb:
                        color = f"#{rgb}"
                        if color not in colors:
                            colors.append(color)
                except Exception:
                    pass
                shape_rows.append({
                    "type": stype,
                    "text_preview": text[:160],
                    "left": int(getattr(shape, "left", 0) or 0),
                    "top": int(getattr(shape, "top", 0) or 0),
                    "width": int(getattr(shape, "width", 0) or 0),
                    "height": int(getattr(shape, "height", 0) or 0),
                    "has_chart": bool(getattr(shape, "has_chart", False)),
                    "has_table": bool(getattr(shape, "has_table", False)),
                })
            title = self._slide_title(slide, texts, idx)
            data_key = self._infer_data_key(title + " " + " ".join(texts[:4]))
            slides.append({
                "slide_index": idx,
                "title": title,
                "data_key": data_key,
                "chart_count": chart_count,
                "table_count": table_count,
                "shape_count": len(shape_rows),
                "text_samples": texts[:8],
                "shapes": shape_rows[:80],
            })

        return {
            "template_name": template_name or os.path.basename(pptx_path),
            "reference_pptx_path": pptx_path,
            "analyzed_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "slide_count": len(slides),
            "chart_count": sum(int(s.get("chart_count") or 0) for s in slides),
            "table_count": sum(int(s.get("table_count") or 0) for s in slides),
            "color_scheme": {"colors": colors[:20]},
            "slides": slides,
        }

    def generate_pptx(self, reference_pptx_path: str, preview_payload: Dict[str, Any], output_dir: str, label: str = "") -> Dict[str, Any]:
        """Generate by preserving the reference deck and editing only detected data.

        Important: this must NOT create a new style or append artificial footer
        boxes. The requirement is "same as reference PPT"; therefore all slides,
        shapes, themes, images, charts, placeholders, transitions, and manual
        text stay as-is unless the text looks like a meta/build/crash/MTBF field
        that should be refreshed for the selected build data.
        """
        if not os.path.exists(reference_pptx_path):
            raise FileNotFoundError(f"Reference PPTX not found: {reference_pptx_path}")

        os.makedirs(output_dir, exist_ok=True)
        target = _safe_str(preview_payload.get("target")) or "target"
        metas = [m for m in (preview_payload.get("selected_metas") or []) if isinstance(m, dict)]
        meta_label = "_".join(_safe_name(_safe_str(m.get("meta_id"))) for m in metas[:6]) or "coredeck"
        if len(metas) > 6:
            meta_label += f"_plus{len(metas) - 6}"
        stamp = datetime.now().strftime("%Y-%m-%d_%H%M%S")
        out_name = f"{_safe_name(target)}_{_safe_name(label or meta_label)}_{stamp}.pptx"
        output_path = os.path.join(output_dir, out_name)

        # Do NOT open/save through python-pptx here. python-pptx can rewrite the
        # package and subtly change layout/chart/theme internals. For exact
        # reference matching, copy the PPTX byte-for-byte first, then patch only
        # slide XML text nodes (`a:t`) that contain data values.
        shutil.copy2(reference_pptx_path, output_path)
        replacements = self._build_reference_replacements(preview_payload)
        changed_shapes = self._patch_pptx_text_xml_in_place(output_path, replacements)
        changed_tables = 0
        return {
            "output_path": output_path,
            "file_name": out_name,
            "changed_shapes": changed_shapes,
            "changed_tables": changed_tables,
            "mode": "reference-preserve-update-data-only",
        }

    def _slide_title(self, slide: Any, texts: List[str], idx: int) -> str:
        try:
            if slide.shapes.title and _safe_str(slide.shapes.title.text):
                return _safe_str(slide.shapes.title.text)
        except Exception:
            pass
        for text in texts:
            first = _safe_str(text).splitlines()[0] if _safe_str(text) else ""
            if first:
                return first[:120]
        return f"Slide {idx}"

    def _infer_data_key(self, text: str) -> str:
        low = _safe_str(text).lower()
        best = ("generic", 0)
        for key, words in self.DATA_KEYWORDS.items():
            score = sum(1 for w in words if w in low)
            if score > best[1]:
                best = (key, score)
        return best[0]

    def _build_reference_replacements(self, data: Dict[str, Any]) -> Dict[str, Any]:
        metas = [m for m in (data.get("selected_metas") or []) if isinstance(m, dict)]
        counts = data.get("summary_counts") or {}
        first = metas[0] if metas else {}
        meta_ids = [_safe_str(m.get("meta_id")) for m in metas if _safe_str(m.get("meta_id"))]
        build_ids = []
        for m in metas:
            for b in m.get("build_ids") or []:
                b = _safe_str(b)
                if b and b not in build_ids:
                    build_ids.append(b)
        total_crashes = sum(int(m.get("crashes") or 0) for m in metas)
        total_devices = sum(int(m.get("device_count") or 0) for m in metas)
        mtbf_values = [_safe_str(m.get("mtbf")) for m in metas if _safe_str(m.get("mtbf"))]
        return {
            "target": _safe_str(data.get("target")),
            "generated_at": _safe_str(data.get("generated_at")),
            "meta_ids": meta_ids,
            "build_ids": build_ids,
            "first_meta": _safe_str(first.get("meta_id")),
            "first_build": _safe_str((first.get("build_ids") or [""])[0] if first else ""),
            "total_crashes": total_crashes,
            "total_devices": total_devices,
            "first_crashes": int(first.get("crashes") or 0) if first else 0,
            "first_devices": int(first.get("device_count") or 0) if first else 0,
            "first_mtbf": _safe_str(first.get("mtbf")) or "NA",
            "mtbf_text": ", ".join(mtbf_values) if mtbf_values else "NA",
            "total_jiras": counts.get("total_jiras", 0),
            "open_jiras": counts.get("open_jiras", 0),
            "unique_crs": counts.get("unique_crs", 0),
        }

    def _patch_pptx_text_xml_in_place(self, pptx_path: str, replacements: Dict[str, Any]) -> int:
        changed_nodes = 0
        tmp_path = pptx_path + ".tmp"
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if re.match(r"ppt/slides/slide\d+\.xml$", item.filename):
                    try:
                        root = ET.fromstring(data)
                        changed = False
                        for node in root.findall(".//a:t", ns):
                            original = node.text or ""
                            updated = self._replace_reference_text(original, replacements)
                            if updated != original:
                                node.text = updated
                                changed = True
                                changed_nodes += 1
                        if changed:
                            data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                    except Exception:
                        pass
                zout.writestr(item, data)
        os.replace(tmp_path, pptx_path)
        return changed_nodes

    def _replace_reference_text(self, text: str, r: Dict[str, Any]) -> str:
        """Replace only reference data tokens; keep all other wording/style.

        Examples preserved:
        - "Meta-123" -> selected first/new meta
        - "Crashes: 12" -> "Crashes: <new>"
        - "MTBF - 50" -> "MTBF - <new>"
        - "Total JIRAs: 100" -> refreshed KPI
        Everything not matching these data patterns is returned unchanged.
        """
        if not text:
            return text
        out = text

        first_meta = _safe_str(r.get("first_meta"))
        if first_meta:
            out = re.sub(r"(?i)\bMeta[-_ ]?\d{1,6}\b", first_meta, out)
            out = re.sub(r"(?i)\bMETA[-_ ]?\d{1,6}\b", first_meta, out)

        first_build = _safe_str(r.get("first_build"))
        if first_build:
            out = re.sub(r"(?i)(\bbuild(?:\s*id)?\s*[:\-]\s*)([^\n\r;|,]+)", lambda m: m.group(1) + first_build, out)

        meta_ids = ", ".join(r.get("meta_ids") or [])
        if meta_ids:
            out = re.sub(r"(?i)(\bmetas?\s*[:\-]\s*)([^\n\r]+)", lambda m: m.group(1) + meta_ids, out)

        build_ids = ", ".join((r.get("build_ids") or [])[:4])
        if build_ids:
            out = re.sub(r"(?i)(\bbuilds?\s*[:\-]\s*)([^\n\r]+)", lambda m: m.group(1) + build_ids, out)

        numeric_patterns = [
            (r"(?i)(\btotal\s+crashes?\s*[:\-]\s*)([\d,]+)", r.get("total_crashes")),
            (r"(?i)(\bcrashes?\s*[:\-]\s*)([\d,]+)", r.get("first_crashes")),
            (r"(?i)(\bcrash\s+count\s*[:\-]\s*)([\d,]+)", r.get("first_crashes")),
            (r"(?i)(\bdevices?\s*[:\-]\s*)([\d,]+)", r.get("first_devices")),
            (r"(?i)(\bdevice\s+count\s*[:\-]\s*)([\d,]+)", r.get("first_devices")),
            (r"(?i)(\btotal\s+jiras?\s*[:\-]\s*)([\d,]+)", r.get("total_jiras")),
            (r"(?i)(\bopen\s+jiras?\s*[:\-]\s*)([\d,]+)", r.get("open_jiras")),
            (r"(?i)(\bunique\s+crs?\s*[:\-]\s*)([\d,]+)", r.get("unique_crs")),
        ]
        for pattern, value in numeric_patterns:
            if value is not None:
                out = re.sub(pattern, lambda m, v=value: m.group(1) + str(v), out)

        mtbf = _safe_str(r.get("first_mtbf"))
        if mtbf:
            out = re.sub(r"(?i)(\bmtbf\s*[:\-]\s*)([<>\d.,NAna/ ]+)", lambda m: m.group(1) + mtbf, out)

        return out
